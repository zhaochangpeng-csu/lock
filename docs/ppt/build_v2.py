#!/usr/bin/env python3
"""
Build the 46-page project deck ON TOP of the template "14 (1).pptx".

Key difference from v1: instead of wiping the template and drawing slides from
scratch (which destroyed every layout and animation), every slide here is a
faithful duplicate of a template page - background, decorative shapes, groups
and <p:timing> animations are all preserved. Only text is replaced, and dense
content (tables / long lists) is placed in a panel that matches the template's
light-blue styling.
"""
import copy
import json
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# Template palette (from the theme of 14 (1).pptx)
C_DK = RGBColor(0x44, 0x54, 0x6A)
C_LT = RGBColor(0xFF, 0xFF, 0xFF)
C_A1 = RGBColor(0xAD, 0xC7, 0xDD)
C_A2 = RGBColor(0xE1, 0xF1, 0xFE)
C_A4 = RGBColor(0x6A, 0x9A, 0xC4)
C_A5 = RGBColor(0xA3, 0xC9, 0xE9)

# Prototype pages (0-based index into the template's slide list)
P_COVER, P_AGENDA, P_SECTION = 0, 1, 2
P_LIST3 = 5      # p06  3 content text boxes
P_COLS4 = 8      # p09  4 columns
P_ROWS4 = 10     # p11  4 bottom groups
P_2COL = 11      # p12  left frame + 2 (title, text) pairs
P_IMGR = 13      # p14  Mockup 14
P_IMGL = 14      # p15  Real Estate 5 - big picture placeholder
P_GRID4 = 15     # p16  title + 4 groups
P_CHART5 = 18    # p19  chart + 5 legend rows
P_FRAME = 21     # p22  big framed area (best for tables / long lists)
P_END = 22       # p23


# ── low level helpers ────────────────────────────────────
# Relationship types that must NOT be carried over to a duplicated slide.
# Copying themeOverride makes PowerPoint declare the file corrupt; notesSlide
# would drag the template's speaker notes along; slideLayout is already bound
# by add_slide().
_SKIP_REL_SUFFIX = ("/themeOverride", "/notesSlide", "/slideLayout")


def _copy_rels(dest_slide, src_slide):
    """Re-create the source slide's relationships on the duplicate.

    Iterates the source relationships (not just the r:id attributes) so that
    non-referenced ones such as `tags` come along too - omitting tags has been
    observed to make PowerPoint reject the file.
    """
    mapping = {}
    for rid, rel in list(src_slide.part.rels.items()):
        if any(rel.reltype.endswith(sfx) for sfx in _SKIP_REL_SUFFIX):
            continue
        if rel.is_external:
            continue
        try:
            new_rid = dest_slide.part.relate_to(rel.target_part, rel.reltype)
        except Exception:
            continue
        mapping[rid] = new_rid

    # Remap every relationship reference in the copied slide XML.
    if mapping:
        for el in dest_slide._element.iter():
            for attr, val in list(el.attrib.items()):
                if val and attr.startswith("{%s}" % R_NS) and val in mapping:
                    el.set(attr, mapping[val])
    return mapping


def _next_ids(slide, count):
    """Return `count` fresh shape ids not used anywhere on the slide."""
    used = set()
    for el in slide._element.iter():
        if el.tag == qn("p:cNvPr"):
            try:
                used.add(int(el.get("id")))
            except (TypeError, ValueError):
                pass
    nxt = max(used) + 1 if used else 1000
    return list(range(nxt, nxt + count))


def duplicate(prs, src_slide):
    dest = prs.slides.add_slide(src_slide.slide_layout)
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    src_cSld = src_slide._element.find(qn("p:cSld"))
    dst_cSld = dest._element.find(qn("p:cSld"))

    for bg in dst_cSld.findall(qn("p:bg")):
        dst_cSld.remove(bg)
    src_bg = src_cSld.find(qn("p:bg"))
    if src_bg is not None:
        dst_cSld.insert(0, copy.deepcopy(src_bg))

    src_tree = src_cSld.find(qn("p:spTree"))
    dst_tree = dst_cSld.find(qn("p:spTree"))
    for child in list(dst_tree):
        dst_tree.remove(child)
    for child in src_tree:
        dst_tree.append(copy.deepcopy(child))

    dst_timing = dest._element.find(qn("p:timing"))
    if dst_timing is not None:
        dest._element.remove(dst_timing)
    src_timing = src_slide._element.find(qn("p:timing"))
    if src_timing is not None:
        dest._element.append(copy.deepcopy(src_timing))

    _copy_rels(dest, src_slide)
    return dest


def find_shape(slide, sid):
    def rec(shapes):
        for sh in shapes:
            if sh.shape_id == sid:
                return sh
            if str(sh.shape_type).startswith("GROUP"):
                r = rec(sh.shapes)
                if r is not None:
                    return r
        return None
    return rec(slide.shapes)


def _proto_run(shape):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            return r
    return None


def set_text(shape, lines, size=None, bold=None, color=None, align=None, width=None, spacing=None):
    """Replace text but keep the template's run formatting (font, colour, spacing)."""
    if shape is None:
        return
    if isinstance(lines, str):
        lines = [lines]
    tf = shape.text_frame
    proto = _proto_run(shape)
    orig_align = tf.paragraphs[0].alignment if tf.paragraphs else None
    tf.clear()
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = align if align is not None else (orig_align or PP_ALIGN.LEFT)
        if spacing:
            p.space_after = Pt(spacing)
        for r in p.runs:
            if proto is not None:
                src = proto._r.find(qn("a:rPr"))
                if src is not None:
                    old = r._r.find(qn("a:rPr"))
                    if old is not None:
                        r._r.remove(old)
                    r._r.insert(0, copy.deepcopy(src))
            if size:
                r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            if color is not None:
                r.font.color.rgb = color
    if width:
        shape.width = Inches(width)


def clear_text(shape):
    if shape is not None:
        shape.text_frame.clear()


def all_text_shapes(shapes):
    """Every text-bearing shape inside a group, at any nesting depth."""
    out = []
    for sh in shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            out.append(sh)
        if str(sh.shape_type).startswith("GROUP"):
            out += all_text_shapes(sh.shapes)
    return out


def add_textbox(slide, left, top, width, height, lines, size=11, color=C_DK,
                bold=False, align=PP_ALIGN.LEFT, spacing=3, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = "DengXian"
            r.font.color.rgb = color
    return tb


def add_panel(slide, left, top, width, height, fill=C_LT, line=C_A1, radius=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def add_page_no(slide, n):
    tb = slide.shapes.add_textbox(Inches(9.05), Inches(5.28), Inches(0.6), Inches(0.22))
    p = tb.text_frame.paragraphs[0]
    p.text = str(n)
    p.alignment = PP_ALIGN.RIGHT
    for r in p.runs:
        r.font.size = Pt(8)
        r.font.name = "DengXian"
        r.font.color.rgb = C_A4


# ── deck builder ─────────────────────────────────────────
class Builder:
    def __init__(self, tpl_path):
        self.prs = Presentation(tpl_path)
        self.protos = list(self.prs.slides)
        # remember the template's demo pages so they can be stripped at the end
        self._orig_rids = [sid.rId for sid in self.prs.slides._sldIdLst]
        self.out = []
        self.max_id = 5000

    def _dup(self, idx):
        s = duplicate(self.prs, self.protos[idx])
        self.out.append(s)
        return s

    # ---- cover / ending -------------------------------------------------
    def cover(self, title, sub_cn, sub_en, author, date):
        s = self._dup(P_COVER)
        set_text(find_shape(s, 50), title)
        set_text(find_shape(s, 65), sub_cn)
        set_text(find_shape(s, 51), sub_en)
        set_text(find_shape(s, 13), "汇报人")
        set_text(find_shape(s, 14), f"{author}   {date}")
        return s

    def ending(self, title, sub_cn, sub_en, author):
        s = self._dup(P_END)
        set_text(find_shape(s, 50), title)
        set_text(find_shape(s, 65), sub_cn)
        set_text(find_shape(s, 51), sub_en)
        set_text(find_shape(s, 13), "汇报人")
        set_text(find_shape(s, 14), author)
        return s

    # ---- agenda ---------------------------------------------------------
    def agenda(self, topics):
        s = self._dup(P_AGENDA)
        set_text(find_shape(s, 16), "目录")
        set_text(find_shape(s, 17), "CONTENTS")
        groups = [20, 22, 23, 26]
        # Clone enough groups to hold every topic.
        src_tree = s.shapes._spTree
        base = None
        for sh in s.shapes:
            if sh.shape_id == 20:
                base = sh
                break
        while len(groups) < len(topics):
            new_el = copy.deepcopy(base._element)
            # Renumber every cNvPr in the clone so ids stay unique. The group's
            # own cNvPr is the first one encountered - that is the id we track.
            first_id = None
            for el in new_el.iter():
                if el.tag == qn("p:cNvPr"):
                    self.max_id += 1
                    el.set("id", str(self.max_id))
                    if first_id is None:
                        first_id = self.max_id
            src_tree.append(new_el)
            if first_id is not None:
                groups.append(first_id)
            else:
                break
        # Layout: 7 rows evenly spread on the right half
        y0, step = 0.72, 0.63
        for i, gid in enumerate(groups[:len(topics)]):
            g = find_shape(s, gid)
            if g is None:
                continue
            g.top = Inches(y0 + i * step)
            g.left = Inches(4.90)
        # Fill text: each group holds (number, title, desc)
        pairs = {20: (57, 59, 60), 22: (51, 53, 54), 23: (39, 44, 47), 26: (27, 33, 34)}
        for i, gid in enumerate(groups[:len(topics)]):
            g = find_shape(s, gid)
            if g is None:
                continue
            # Inside each agenda group the order is: number / title / description
            texts = [sh for sh in g.shapes if sh.has_text_frame]
            if len(texts) >= 3:
                num_t, title_t, desc_t = texts[0], texts[1], texts[2]
            if num_t is not None:
                set_text(num_t, f"{i+1:02d}")
            if title_t is not None:
                set_text(title_t, topics[i])
            if desc_t is not None:
                clear_text(desc_t)
        return s

    # ---- section --------------------------------------------------------
    def section(self, num, title, subtitle):
        s = self._dup(P_SECTION)
        # The template renders a single oversized digit at 215pt; two digits
        # would be ~3.3in wide and break the composition, so keep one digit.
        set_text(find_shape(s, 2), str(num))
        set_text(find_shape(s, 3), str(num))
        set_text(find_shape(s, 25), title)
        set_text(find_shape(s, 27), subtitle)
        return s

    # ---- 3-item list (p06) ---------------------------------------------
    def list3(self, title, blocks):
        """blocks: list of (heading, [lines]) - at most 3."""
        s = self._dup(P_LIST3)
        ids = [29, 30, 31]
        for i, sid in enumerate(ids):
            sh = find_shape(s, sid)
            if sh is None:
                continue
            if i < len(blocks):
                head, lines = blocks[i]
                set_text(sh, [head] + list(lines) if lines else [head])
            else:
                clear_text(sh)
        # p06's header block (group 34) holds the template's "输入标题" stub.
        clear_text(find_shape(s, 37))
        clear_text(find_shape(s, 38))
        # Title goes below the decorative header (ends y=0.78) and clear of the
        # left image group.
        add_textbox(s, 0.70, 0.84, 4.30, 0.32, [title], size=14, bold=True, color=C_A4)
        return s

    # ---- 4 columns (p09) ------------------------------------------------
    def cols4(self, title, blocks):
        s = self._dup(P_COLS4)
        # p09 exposes a full-width title band (id 87) and a subtitle bar (id 89).
        set_text(find_shape(s, 87), title, size=14, bold=True, color=C_A4)
        clear_text(find_shape(s, 89))
        clear_text(find_shape(s, 107))
        # header block stubs
        clear_text(find_shape(s, 15))
        clear_text(find_shape(s, 16))
        ids = [112, 119, 126, 132]
        for i, sid in enumerate(ids):
            g = find_shape(s, sid)
            if g is None:
                continue
            # Text boxes may sit at different nesting depths per column, and the
            # body box precedes the heading, so match on the template's stub copy.
            inner = all_text_shapes(g.shapes)
            header = next((x for x in inner if "Header here please" in x.text_frame.text), None)
            bodies = [x for x in inner if "We are placing" in x.text_frame.text]
            if i < len(blocks):
                head, lines = blocks[i]
                if header is not None:
                    set_text(header, head, size=12, bold=True, color=C_A4)
                if bodies:
                    set_text(bodies[0], lines, size=9, color=C_DK, spacing=2)
                    for extra in bodies[1:]:
                        clear_text(extra)
            else:
                for t in inner:
                    clear_text(t)
        return s

    # ---- 2 columns (p12) ------------------------------------------------
    def twocol(self, title, left, right):
        s = self._dup(P_2COL)
        # header block stubs + title placed left of / below the header block
        clear_text(find_shape(s, 17))
        clear_text(find_shape(s, 19))
        add_textbox(s, 0.85, 0.84, 4.30, 0.32, [title], size=14, bold=True, color=C_A4)
        # p12: (25 title, 24 body) and (5 title, 4 body)
        set_text(find_shape(s, 25), left.get("head", ""), size=12, bold=True, color=C_A4)
        set_text(find_shape(s, 5), right.get("head", ""), size=12, bold=True, color=C_A4)
        b1 = find_shape(s, 24)
        b2 = find_shape(s, 4)
        for sh, items, top in ((b1, left.get("items", []), 1.66), (b2, right.get("items", []), 3.28)):
            if sh is None:
                continue
            sh.top = Inches(top)
            sh.height = Inches(1.42 if len(items) > 5 else 1.00)
            set_text(sh, items, size=9, color=C_DK, spacing=2)
        return s

    # ---- framed panel page (p22) ---------------------------------------
    def frame(self, title, blocks=None, table=None, note=None, timeline=None,
              conclusion=None, bullets=None):
        s = self._dup(P_FRAME)
        # header block stubs (the title itself lives inside the panel, which is
        # wide enough for long technical titles)
        clear_text(find_shape(s, 14))
        clear_text(find_shape(s, 17))
        # p22's big framed area is (1.04, 0.95) 5.77 x 4.13. The title lives
        # inside the panel so it never collides with the decorative header.
        x, y, w, h = 1.04, 0.95, 5.77, 4.13
        add_panel(s, x + 0.09, y + 0.09, w - 0.18, h - 0.18)
        add_textbox(s, x + 0.26, y + 0.20, w - 0.52, 0.36, [title],
                    size=15, bold=True, color=C_A4)
        ix, iy, iw, ih = x + 0.26, y + 0.62, w - 0.52, h - 0.80

        if table:
            headers, rows = table
            self._table(s, headers, rows, ix, iy, iw, ih)
        elif timeline:
            self._timeline(s, timeline, ix, iy, iw, ih)
        elif bullets is not None:
            self._bullets(s, bullets, ix, iy, iw, ih, conclusion)
        elif blocks:
            self._blocks(s, blocks, ix, iy, iw, ih)

        # right-hand column of p22
        set_text(find_shape(s, 129), note or "", size=9, color=C_DK)
        return s

    def _table(self, s, headers, rows, x, y, w, h):
        nrows, ncols = len(rows) + 1, len(headers)
        gf = s.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(min(h, 0.34 * nrows)))
        tbl = gf.table
        try:
            tbl.first_row = True
            tbl.horz_banding = True
        except Exception:
            pass
        for i, hd in enumerate(headers):
            c = tbl.cell(0, i)
            c.text = str(hd)
            c.fill.solid()
            c.fill.fore_color.rgb = C_A4
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(8.5)
                p.font.bold = True
                p.font.color.rgb = C_LT
                p.alignment = PP_ALIGN.CENTER
        for ri, row in enumerate(rows, 1):
            for ci, val in enumerate(row):
                c = tbl.cell(ri, ci)
                c.text = str(val)
                c.margin_left = Inches(0.04)
                c.margin_right = Inches(0.04)
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.font.name = "DengXian"
                    p.font.color.rgb = C_DK
                    p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
        return tbl

    def _blocks(self, s, blocks, x, y, w, h):
        cur = y
        for head, lines in blocks:
            add_textbox(s, x, cur, w, 0.24, [head], size=11, bold=True, color=C_A4, spacing=1)
            cur += 0.26
            for ln in lines:
                add_textbox(s, x + 0.16, cur, w - 0.16, 0.22, [ln], size=9, color=C_DK, spacing=1)
                cur += 0.215
            cur += 0.07
        return cur

    def _bullets(self, s, items, x, y, w, h, conclusion=None):
        cur = y
        box_h = h - 0.75 if conclusion else h
        for it in items:
            txt = it if isinstance(it, str) else it.get("t", str(it))
            add_textbox(s, x, cur, w, 0.24, ["· " + txt], size=10.5, color=C_DK, spacing=1)
            cur += 0.30
        if conclusion:
            add_panel(s, x, y + box_h + 0.06, w, 0.62, fill=C_A4, line=None)
            add_textbox(s, x + 0.12, y + box_h + 0.14, w - 0.24, 0.5, [conclusion],
                        size=10.5, bold=True, color=C_LT, align=PP_ALIGN.CENTER, spacing=0)

    def _timeline(self, s, milestones, x, y, w, h):
        n = len(milestones)
        if not n:
            return
        line_y = y + h * 0.52
        add_panel(s, x, line_y - 0.015, w, 0.03, fill=C_A1, line=None, radius=False)
        step = w / n
        for i, ms in enumerate(milestones):
            cx = x + step * i + step / 2
            label = ms.get("label", "")
            desc = ms.get("desc", "")
            active = ms.get("active", False)
            col = C_A4 if active else C_A1
            dot = add_panel(s, cx - 0.09, line_y - 0.09, 0.18, 0.18, fill=col, line=None, radius=False)
            dot.shape_type  # touch to keep reference
            add_textbox(s, cx - step / 2 + 0.04, line_y - 0.62, step - 0.08, 0.3,
                        [label], size=11, bold=True, color=C_DK, align=PP_ALIGN.CENTER, spacing=1)
            if desc:
                add_textbox(s, cx - step / 2 + 0.04, line_y + 0.22, step - 0.08, 0.42,
                            [desc], size=8.5, color=C_DK, align=PP_ALIGN.CENTER, spacing=1)

    # ---- chart page (p19) ----------------------------------------------
    def chart(self, title, chart_type, categories, series, legend=None):
        s = self._dup(P_CHART5)
        clear_text(find_shape(s, 31))   # header block title stub
        clear_text(find_shape(s, 32))   # header block subtitle stub
        add_textbox(s, 0.85, 0.84, 4.70, 0.34, [title], size=15, bold=True, color=C_A4)
        # drop the template's bubble chart, put ours in the same spot
        old = None
        for sh in s.shapes:
            if getattr(sh, "has_chart", False) and sh.has_chart:
                old = sh
                break
        box = (0.82, 1.58, 4.26, 2.70)
        if old is not None:
            box = (old.left / 914400.0, old.top / 914400.0,
                   old.width / 914400.0, old.height / 914400.0)
            old._element.getparent().remove(old._element)
        cd = CategoryChartData()
        cd.categories = categories
        for se in series:
            cd.add_series(se.get("name", ""), se.get("values", []))
        ct = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
              "line": XL_CHART_TYPE.LINE_MARKERS,
              "pie": XL_CHART_TYPE.PIE}.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        gf = s.shapes.add_chart(ct, Inches(box[0]), Inches(box[1]), Inches(box[2]), Inches(box[3]), cd)
        ch = gf.chart
        ch.has_legend = chart_type != "pie"
        if ch.has_legend:
            ch.legend.position = XL_LEGEND_POSITION.BOTTOM
            ch.legend.include_in_layout = False
        try:
            ch.font.size = Pt(9)
        except Exception:
            pass
        if chart_type == "pie":
            plot = ch.plots[0]
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.show_percentage = True
            dl.show_value = False
            try:
                dl.font.size = Pt(9)
            except Exception:
                pass
        # right-hand legend rows of p19 (groups 11,14,17,20,23)
        if legend:
            for i, gid in enumerate([11, 14, 17, 20, 23]):
                g = find_shape(s, gid)
                if g is None:
                    continue
                texts = [sh for sh in g.shapes if sh.has_text_frame]
                if i < len(legend) and texts:
                    set_text(texts[0], legend[i], size=9, color=C_DK)
                else:
                    for t in texts:
                        clear_text(t)
        return s

    # ---- image page (p15) ----------------------------------------------
    def image(self, title, path, caption=None):
        s = self._dup(P_IMGL)
        # Clear the template's demo copy (native title slot + lorem ipsum).
        for sid in (33, 34, 10, 17, 26, 27, 28):
            clear_text(find_shape(s, sid))
        # Title sits below the decorative header block (which ends at y=0.78).
        add_textbox(s, 0.72, 0.84, 4.70, 0.38, [title], size=15, bold=True, color=C_A4)
        ph = None
        for sh in s.shapes:
            if sh.is_placeholder and sh.placeholder_format.type is not None and \
               "PICTURE" in str(sh.placeholder_format.type):
                ph = sh
                break
        if ph is not None:
            px, py = ph.left / 914400.0, ph.top / 914400.0
            pw, phh = ph.width / 914400.0, ph.height / 914400.0
            keep_id = ph.shape_id
            # Fit the image inside the placeholder box (no stretching).
            dw, dh = pw, phh
            try:
                from PIL import Image as PILImage
                with PILImage.open(path) as im:
                    iw, ih = im.size
                sc = min(pw / iw, phh / ih)
                dw, dh = iw * sc, ih * sc
            except Exception:
                pass
            cx, cy = px + (pw - dw) / 2, py + (phh - dh) / 2
            ph._element.getparent().remove(ph._element)
            pic = s.shapes.add_picture(path, Inches(cx), Inches(cy), Inches(dw), Inches(dh))
            # Re-use the placeholder's shape id so the template's animation
            # (which targets the shape by id) keeps pointing at this picture.
            for el in pic._element.iter():
                if el.tag == qn("p:cNvPr"):
                    el.set("id", str(keep_id))
                    break
        if caption:
            set_text(find_shape(s, 30), caption, size=8.5, color=C_DK)
        return s

    def save(self, path):
        """Strip the template's own demo pages, leaving only the built deck."""
        sldIdLst = self.prs.slides._sldIdLst
        for rid in self._orig_rids:
            for sid in list(sldIdLst):
                if sid.rId == rid:
                    try:
                        self.prs.part.drop_rel(rid)
                    except Exception:
                        pass
                    sldIdLst.remove(sid)
                    break
        self.prs.save(path)
        return len(sldIdLst)


def _blocks(items):
    return [(it.get("t", ""), it.get("s", [])) for it in items]


def main():
    src = os.path.join(ROOT, "docs", "ppt", "slides.json")
    out = os.path.join(ROOT, "docs", "智能门锁项目_立项与技术汇报_模板版.pptx")
    if len(sys.argv) > 1:
        src = sys.argv[1]
    if len(sys.argv) > 2:
        out = sys.argv[2]

    data = json.load(open(src, encoding="utf-8"))
    b = Builder(os.path.join(ROOT, "14 (1).pptx"))
    chapter = ""

    for s in data.get("slides", []):
        typ = s.get("type", "content")
        title = s.get("title", "")

        if typ == "cover":
            sub = s.get("subtitle", "").split("\n")
            b.cover(title,
                    sub[0] if len(sub) > 0 else "",
                    sub[1] if len(sub) > 1 else "",
                    s.get("author", ""), s.get("date", ""))

        elif typ == "agenda":
            b.agenda(s.get("topics", []))

        elif typ == "section":
            # "01  立项背景与市场机会" -> number 1, name for the side rail
            head, _, name = title.partition("  ")
            try:
                num = int("".join(ch for ch in head if ch.isdigit()) or 0)
            except ValueError:
                num = 0
            chapter = name.strip() or title
            b.section(num, title, s.get("subtitle", ""))

        elif typ == "two_column":
            b.twocol(title,
                     {"head": s.get("left_head", ""), "items": s.get("left", [])},
                     {"head": s.get("right_head", ""), "items": s.get("right", [])})

        elif typ == "table":
            b.frame(title, table=(s.get("headers", []), s.get("rows", [])), note=chapter)

        elif typ == "chart":
            b.chart(title, s.get("chart_type", "bar"), s.get("categories", []),
                    s.get("series", []), legend=s.get("legend"))

        elif typ == "image":
            b.image(title, s.get("image_path", ""), s.get("caption"))

        elif typ == "timeline":
            b.frame(title, timeline=s.get("milestones", []), note=chapter)

        elif typ == "summary":
            b.frame(title, bullets=s.get("points", []),
                    conclusion=s.get("conclusion", ""), note=chapter)

        elif typ == "contact":
            info = s.get("info", [])
            b.ending(title,
                     info[0] if info else "",
                     " / ".join(info[1:3]) if len(info) > 1 else "",
                     info[3] if len(info) > 3 else "")

        else:  # content
            items = s.get("items", [])
            if items and isinstance(items[0], dict):
                blocks = _blocks(items)
                if len(blocks) <= 3:
                    b.list3(title, blocks)
                elif len(blocks) == 4:
                    b.cols4(title, blocks)
                else:
                    b.frame(title, blocks=blocks, note=chapter)
            else:
                if len(items) <= 3:
                    b.list3(title, [(str(i), []) for i in items])
                elif len(items) == 4:
                    b.cols4(title, [(str(i), []) for i in items])
                else:
                    b.frame(title, bullets=[str(i) for i in items], note=chapter)

    for i, s in enumerate(b.out, 1):
        add_page_no(s, i)

    n = b.save(out)
    print(f"已生成 {out}  ({n} 页)")


if __name__ == "__main__":
    main()
