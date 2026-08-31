#!/usr/bin/env python3
"""Layout sanity check for the generated deck.

Detects (a) shapes off the canvas, (b) text boxes whose estimated rendered
height exceeds the space available, and (c) tables that would expand past the
slide bottom once PowerPoint auto-fits wrapped cell text.

Approximate by design - a smoke test, not a renderer.
"""
import sys
from pptx import Presentation

EMU_IN = 914400.0
SLIDE_BOTTOM_PAD = 0.15


def char_em(ch):
    # CJK / full-width punctuation occupies ~1em, latin ~0.55em.
    return 1.0 if ord(ch) > 0x2E80 else 0.55


def wrapped_lines(text, pt, width_in):
    em_in = pt / 72.0
    width_em = sum(char_em(c) for c in text)
    if width_em <= 0:
        return 1
    return max(1, -(-int(width_em * em_in * 1000) // int(max(width_in, 0.3) * 1000)))


def est_text_height(tf, box_w_in, default_pt=18.0):
    total = 0.0
    for p in tf.paragraphs:
        text = p.text or ""
        if not text.strip():
            total += default_pt * 1.30
            continue
        # Size may be set on the run OR inherited from the paragraph defRPr.
        sizes = [r.font.size.pt for r in p.runs if r.font.size]
        if not sizes and p.font.size:
            sizes = [p.font.size.pt]
        pt = max(sizes) if sizes else default_pt
        indent = 0.25 * p.level
        avail = max(box_w_in - indent, 0.3)
        lines = wrapped_lines(text, pt, avail)
        sa = p.space_after.pt if p.space_after else 0
        total += lines * pt * 1.28 + sa
    return total / 72.0


def est_table_height(table):
    """Estimate height after PowerPoint auto-grows rows to fit wrapped text."""
    n_cols = len(table.columns)
    col_w = [c.width / EMU_IN for c in table.columns]
    total = 0.0
    for row in table.rows:
        max_h = 0.25  # minimum row height
        for ci, cell in enumerate(row.cells):
            tf = cell.text_frame
            text = tf.text or ""
            sizes = [r.font.size.pt for p in tf.paragraphs for r in p.runs if r.font.size]
            sizes += [p.font.size.pt for p in tf.paragraphs if p.font.size]
            pt = max(sizes) if sizes else 12.0
            w = col_w[ci] if ci < n_cols else 1.0
            inner_w = max(w - 0.16, 0.3)  # cell left/right margins
            lines = wrapped_lines(text, pt, inner_w)
            max_h = max(max_h, lines * pt * 1.25 / 72.0 + 0.1)
        total += max_h
    return total


def main(path):
    prs = Presentation(path)
    sw, sh = prs.slide_width / EMU_IN, prs.slide_height / EMU_IN
    print(f"canvas {sw:.2f} x {sh:.2f} in   slides={len(prs.slides.__iter__.__self__._sldIdLst)}\n")

    problems, tight = [], []
    for idx, slide in enumerate(prs.slides, 1):
        title = ""
        for s in slide.shapes:
            if s.has_text_frame and s.text_frame.text.strip():
                title = s.text_frame.text.strip().split("\n")[0][:36]
                break

        worst = 0.0
        for s in slide.shapes:
            if s.left is None:
                continue
            L, T = s.left / EMU_IN, s.top / EMU_IN
            W, H = s.width / EMU_IN, s.height / EMU_IN

            # Templates deliberately bleed full-page artwork past the canvas
            # edge, so only flag a real overhang on non-bleed shapes.
            overhang = max(-L, -T, (L + W) - sw, (T + H) - sh)
            is_bleed = (W >= sw - 0.05) or (H >= sh - 0.05)
            if overhang > 0.15 and not is_bleed:
                blank = s.has_text_frame and not s.text_frame.text.strip()
                if not blank:
                    problems.append(
                        f"[p{idx:02d}] OFF-CANVAS {type(s).__name__} "
                        f"L={L:.2f} T={T:.2f} W={W:.2f} H={H:.2f} ({title})")

            if getattr(s, "has_table", False) and s.has_table:
                need = est_table_height(s.table)
                avail = sh - T - SLIDE_BOTTOM_PAD
                ratio = need / avail
                worst = max(worst, ratio)
                if need > avail:
                    problems.append(
                        f"[p{idx:02d}] TABLE OVERFLOW need~{need:.2f}in avail={avail:.2f}in ({title})")
                elif ratio > 0.8:
                    tight.append(f"[p{idx:02d}] table {ratio:.0%} of available ({title})")
                continue

            if s.has_text_frame and s.text_frame.text.strip():
                # Skip decorative oversized text (chapter numerals, cover
                # titles): the estimator is unreliable above ~40pt and those
                # shapes overflow their box by design.
                sizes = [r.font.size.pt for p in s.text_frame.paragraphs
                         for r in p.runs if r.font.size]
                sizes += [p.font.size.pt for p in s.text_frame.paragraphs if p.font.size]
                if sizes and max(sizes) > 40:
                    continue
                need = est_text_height(s.text_frame, W)
                avail = sh - T - SLIDE_BOTTOM_PAD
                ratio = need / avail
                worst = max(worst, ratio)
                if need > avail:
                    problems.append(
                        f"[p{idx:02d}] TEXT OVERFLOW need~{need:.2f}in avail={avail:.2f}in "
                        f"({title}) :: {s.text_frame.text[:44]!r}")
                elif ratio > 0.8:
                    tight.append(f"[p{idx:02d}] text  {ratio:.0%} of available ({title})")

    for t in tight:
        print("tight  " + t)
    print()
    for p in problems:
        print("ISSUE  " + p)
    print(f"\n{'OK - no overflow detected' if not problems else f'{len(problems)} issue(s)'}")
    return 1 if problems else 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1]))
