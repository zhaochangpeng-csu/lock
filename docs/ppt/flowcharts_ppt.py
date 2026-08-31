#!/usr/bin/env python3
"""Build 5 NATIVE (editable) flowchart slides in the 14(1).pptx light-blue style.

Every box / diamond / arrow / label is a real PowerPoint shape, so text can be
edited, re-coloured and re-positioned directly in PowerPoint (no raster image).
Canvas 13.333 x 7.5 in (standard 16:9) for large, readable type.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
BG = os.path.join(ROOT, "docs", "ppt", "_tpl_bg", "p01_0b91503e.jpg")

# logical canvas = 1920 x 1080 (16:9), mapped onto 13.333 x 7.5 in
SW, SH = 13.333, 7.5
KX = SW / 1920.0          # px -> inches
KY = SH / 1080.0
def P(x, y): return Inches(x * KX), Inches(y * KY)
def F(px): return Pt(px * 0.5)   # SVG px (at 1920 wide) -> PowerPoint pt

# template palette
C_DK = RGBColor(0x44, 0x54, 0x6A)
C_A1 = RGBColor(0xAD, 0xC7, 0xDD)
C_A2 = RGBColor(0xE1, 0xF1, 0xFE)
C_A4 = RGBColor(0x6A, 0x9A, 0xC4)
C_A5 = RGBColor(0xA3, 0xC9, 0xE9)
C_RED = RGBColor(0xE0, 0x62, 0x5C)
C_GREEN = RGBColor(0x4F, 0xAE, 0x7A)
C_GRAY = RGBColor(0xA6, 0xA6, 0xA6)
C_AMBER = RGBColor(0xE5, 0xA6, 0x4F)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CREAM = RGBColor(0xFD, 0xF5, 0xE6)


def _sz(v):
    """Coerce to a Length; tolerate raw numbers and existing Pt/Emu."""
    return v if hasattr(v, "centipoints") else Pt(v)


def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _set_text(shape, lines, size, color, bold=False, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE, font="DengXian", spacing=2):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        for r in p.runs:
            r.font.size = _sz(size)
            r.font.bold = bold
            r.font.name = font
            r.font.color.rgb = color


def add_bg(slide):
    slide.shapes.add_picture(BG, 0, 0, Inches(SW), Inches(SH))


def add_box(slide, x, y, w, h, title, lines=None, *, fill=C_WHITE, border=C_A1,
            accent=None, title_size=15, body_size=12, title_color=C_DK,
            body_color=C_DK, title_bold=True, body_bold=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, *P(x, y), Inches(w * KX), Inches(h * KY))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(2)
    _no_shadow(sh)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    if accent:
        # left accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *P(x, y + 8), Inches(5 * KX), Inches((h - 16) * KY))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        _no_shadow(bar)
        tf.margin_left = Inches(0.14)
    if isinstance(lines, str):
        lines = [lines] if lines else []
    lines = lines or []
    all_lines = [title] + list(lines)
    for i, ln in enumerate(all_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = PP_ALIGN.CENTER if accent else PP_ALIGN.LEFT
        if i == 0:
            p.space_after = Pt(3)
        else:
            p.space_after = Pt(2)
        for r in p.runs:
            r.font.size = _sz(title_size) if i == 0 else _sz(body_size)
            r.font.bold = title_bold if i == 0 else body_bold
            r.font.name = "DengXian"
            r.font.color.rgb = title_color if i == 0 else body_color
    return sh


def add_num(slide, num, x, y, w, h, title, lines=None, *, title_size=15, body_size=12):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, *P(x, y), Inches(w * KX), Inches(h * KY))
    sh.fill.solid(); sh.fill.fore_color.rgb = C_WHITE
    sh.line.color.rgb = C_A1; sh.line.width = Pt(2)
    _no_shadow(sh)
    # number circle
    r = 26
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, *P(x + 28, y + 22), Inches(r * 2 * KX), Inches(r * 2 * KY))
    c.fill.solid(); c.fill.fore_color.rgb = C_A4
    c.line.fill.background(); _no_shadow(c)
    _set_text(c, str(num), F(28), C_WHITE, bold=True)
    # title
    t = slide.shapes.add_textbox(*P(x + 100, y + 26), Inches((w - 110) * KX), Inches(50 * KY))
    _set_text(t, title, title_size, C_DK, bold=True, align=PP_ALIGN.LEFT)
    # body
    if lines:
        b = slide.shapes.add_textbox(*P(x + 34, y + 96), Inches((w - 50) * KX), Inches((h - 100) * KY))
        _set_text(b, lines, body_size, C_DK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    return sh


def add_diamond(slide, cx, cy, w, h, title, lines=None, *, border=C_A4,
                title_size=14, body_size=11):
    half_w, half_h = w / 2, h / 2
    sh = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, *P(cx - half_w, cy - half_h),
                                Inches(w * KX), Inches(h * KY))
    sh.fill.solid(); sh.fill.fore_color.rgb = C_WHITE
    sh.line.color.rgb = border; sh.line.width = Pt(2)
    _no_shadow(sh)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    if isinstance(lines, str):
        lines = [lines]
    lines = lines or []
    for i, ln in enumerate([title] + list(lines)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(2)
        for r in p.runs:
            r.font.size = _sz(title_size) if i == 0 else _sz(body_size)
            r.font.bold = (i == 0)
            r.font.name = "DengXian"
            r.font.color.rgb = C_DK
    return sh


def add_arrow(slide, x1, y1, x2, y2, color=C_A4, dashed=False, width=2.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, *P(x1, y1), *P(x2, y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "none", "w": "med", "len": "med"}))
    return conn


def add_label(slide, x, y, w, text, size=12, color=C_DK, bold=False, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(*P(x, y), Inches(w * KX), Inches(30 * KY))
    _set_text(tb, text, size, color, bold=bold, align=align)
    return tb


def add_header(slide, num, total, title, subtitle):
    # number circle
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, *P(70, 78), Inches(52 * KX), Inches(52 * KY))
    c.fill.solid(); c.fill.fore_color.rgb = C_A4
    c.line.fill.background(); _no_shadow(c)
    _set_text(c, str(num), F(28), C_WHITE, bold=True)
    # title
    t = slide.shapes.add_textbox(*P(140, 86), Inches(1500 * KX), Inches(60 * KY))
    _set_text(t, title, F(46), C_DK, bold=True, align=PP_ALIGN.LEFT)
    # subtitle
    s = slide.shapes.add_textbox(*P(140, 150), Inches(1500 * KX), Inches(40 * KY))
    _set_text(s, subtitle, F(24), C_A4, align=PP_ALIGN.LEFT)
    # accent line
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *P(140, 196), Inches(460 * KX), Inches(3.5 * KY))
    ln.fill.solid(); ln.fill.fore_color.rgb = C_A4
    ln.line.fill.background(); _no_shadow(ln)
    # badge
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, *P(1720, 88), Inches(160 * KX), Inches(54 * KY))
    b.fill.solid(); b.fill.fore_color.rgb = C_A4
    b.line.fill.background(); _no_shadow(b)
    _set_text(b, f"图 {num} / {total}", F(22), C_WHITE)


def add_footer(slide, label):
    t = slide.shapes.add_textbox(*P(80, 1022), Inches(1300 * KX), Inches(30 * KY))
    _set_text(t, label, F(22), C_A4, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════
# 图1 系统启动
# ════════════════════════════════════════════════════════════
def chart_1(slide):
    add_box(slide, 80, 380, 320, 320, "run_smart_lock.sh", ["一键启动脚本", "按顺序拉起右侧组件", "支持 SSH 与开机自启"], accent=C_A4, title_size=16, body_size=13)
    add_arrow(slide, 400, 540, 468, 540)
    # center 2x2
    add_num(slide, 1, 470, 310, 460, 200, "本地工具网关 :8787", ["Agent 工具调用入口", "current_auth_context · request_unlock"])
    add_num(slide, 2, 950, 310, 460, 200, "异常事件服务 :8790", ["接收认证失败事件", "原子写入 latest_event.json"])
    add_num(slide, 3, 470, 560, 460, 200, "语音 Agent 预热", ["预加载 FunASR 模型", "--wait-auth 等待凭证后开麦"])
    add_num(slide, 4, 950, 560, 460, 200, "GUI 启动", ["PySide6 界面·手动/自动认证", "预加载人脸/活体/声纹模型"])
    # bidirectional arrows
    add_arrow(slide, 1410, 525, 1488, 525)
    add_arrow(slide, 1488, 555, 1410, 555)
    add_box(slide, 1490, 380, 350, 320, "进程监督", ["守护四大组件进程", "异常退出自动拉起", "PID 记录 logs/run/..."], fill=C_A2, accent=C_A4, title_size=16, body_size=13)
    # bottom data
    add_label(slide, 80, 828, 300, "数据落盘位置", size=F(24), color=C_DK, bold=True, align=PP_ALIGN.LEFT)
    add_box(slide, 80, 890, 560, 120, "logs/auth_context.json", ["认证凭证 · 每次覆盖为最新一条"], accent=C_A4, title_size=13, body_size=12)
    add_box(slide, 670, 890, 560, 120, "latest_event.json", ["异常事件 · 仅保留最新一条"], accent=C_A4, title_size=13, body_size=12)
    add_box(slide, 1260, 890, 560, 120, "logs/tool_gateway_calls.jsonl", ["工具调用留痕 · 追加式 JSONL"], accent=C_A4, title_size=13, body_size=12)
    add_footer(slide, "启动顺序: 网关 → 事件服务 → Agent 预热 → GUI")


# ════════════════════════════════════════════════════════════
# 图2 待机与触发
# ════════════════════════════════════════════════════════════
def chart_2(slide):
    add_box(slide, 120, 340, 360, 200, "① 红外传感器", ["串口常驻轮询", "9600 波特"], accent=C_A4, title_size=16, body_size=13)
    add_arrow(slide, 480, 440, 560, 440)
    add_box(slide, 560, 340, 360, 200, "② 打开摄像头", ["持续人脸识别", "空闲自动关电源"], accent=C_A4, title_size=16, body_size=13)
    add_arrow(slide, 920, 440, 1000, 440)
    add_diamond(slide, 1180, 440, 360, 260, "红外 / 人脸", ["冷却期过？"])
    add_arrow(slide, 1000, 600, 1000, 720, color=C_RED, dashed=True)
    add_box(slide, 800, 720, 400, 180, "③ 长时间无法识别", ["停留超 3s → 写异常", "离开后重新布防"], border=C_RED, title_size=15, body_size=12)
    add_arrow(slide, 1360, 440, 1520, 440, color=C_GREEN)
    add_box(slide, 1520, 340, 320, 200, "④ 进入五步认证", ["见 图 3"], fill=C_A2, border=C_GREEN, accent=C_GREEN, title_size=16, body_size=13)
    add_arrow(slide, 1000, 570, 800, 570, color=C_GRAY)
    add_label(slide, 860, 548, 100, "否", color=C_GRAY)
    add_arrow(slide, 800, 280, 300, 280, color=C_GRAY)
    add_label(slide, 420, 258, 260, "否 · 继续观察", color=C_GRAY)
    add_box(slide, 1450, 720, 390, 180, "⑤ 其他入口", ["GUI「手动认证」", "main.py --hardware"], border=C_GRAY, title_size=15, body_size=12)
    add_footer(slide, "承接图1 · 启动完成即进入待机轮询")


# ════════════════════════════════════════════════════════════
# 图3 五步认证
# ════════════════════════════════════════════════════════════
def chart_3(slide):
    for i, (num, title) in enumerate([(1, "红外检测"), (2, "人脸识别"), (3, "活体检测"),
                                     (4, "声纹识别"), (5, "融合判定")]):
        x = 100 + i * 350
        add_num(slide, num, x, 330, 320, 180, title, [], title_size=16)
    for x1, x2 in [(160, 190), (510, 540), (860, 890), (1210, 1240)]:
        add_arrow(slide, x1, 510, x2, 510)
    add_arrow(slide, 1500, 600, 1500, 640)
    add_box(slide, 580, 600, 760, 110, "融合得分 = 0.45×人脸 + 0.25×活体 + 0.20×声纹 + 0.10×红外", [], fill=C_A2, border=C_A4, title_size=14, body_size=12)
    add_arrow(slide, 960, 710, 960, 760)
    add_diamond(slide, 960, 850, 420, 200, "融合通过？", ["score ≥ 0.78"])
    add_arrow(slide, 1170, 850, 1320, 850, color=C_GREEN)
    add_label(slide, 1240, 838, 80, "是", color=C_GREEN, bold=True)
    add_box(slide, 1320, 750, 540, 200, "签发一次性认证书", ["credential_id · 300 秒有效"], fill=C_A2, border=C_GREEN, accent=C_GREEN, title_size=16, body_size=13)
    add_arrow(slide, 800, 920, 480, 970, color=C_RED, dashed=True)
    add_label(slide, 660, 930, 80, "否", color=C_RED, bold=True)
    add_box(slide, 100, 960, 380, 80, "认证失败 X", ["写事件 · 上报"], border=C_RED, title_size=14, body_size=12)
    add_footer(slide, "任一环节超时即中止本轮")


# ════════════════════════════════════════════════════════════
# 图4 语音 Agent
# ════════════════════════════════════════════════════════════
def chart_4(slide):
    add_box(slide, 680, 70, 560, 130, "一次性凭证就绪", ["credential_id · 300 秒有效"], fill=C_A2, border=C_GREEN, accent=C_GREEN, title_size=17, body_size=13)
    add_arrow(slide, 960, 200, 960, 280)
    for i, (num, title) in enumerate([(1, "VAD 分段"), (2, "ASR 转写"), (3, "LLM 回复"), (4, "TTS 合成")]):
        x = 100 + i * 440
        add_num(slide, num, x, 280, 400, 160, title, [], title_size=16)
    for x1, x2 in [(500, 540), (940, 980), (1380, 1420)]:
        add_arrow(slide, x1, 360, x2, 360)
    # L-shaped connector from TTS to diamond (approximate with straight)
    add_arrow(slide, 1740, 440, 1740, 480)
    add_arrow(slide, 1740, 480, 960, 480)
    add_arrow(slide, 960, 480, 960, 510)
    add_diamond(slide, 960, 610, 360, 200, "用户意图？", ["说「开门」？"])
    add_arrow(slide, 780, 710, 780, 750, color=C_RED, dashed=True)
    add_arrow(slide, 780, 750, 270, 750, color=C_RED, dashed=True)
    add_label(slide, 800, 735, 120, "凭证无效", color=C_RED)
    add_arrow(slide, 960, 710, 960, 750)
    add_arrow(slide, 960, 750, 740, 750)
    add_label(slide, 830, 735, 120, "说「开门」", color=C_A4)
    add_box(slide, 80, 750, 380, 140, "凭证无效 · 拒绝开门", ["引导重认证 · 通知业主"], border=C_RED, title_size=14, body_size=12)
    add_box(slide, 540, 750, 400, 140, "查询 current_auth_context", ["available · fresh · authorized"], border=C_A4, title_size=14, body_size=12)
    add_box(slide, 1020, 750, 380, 140, "request_unlock", ["→ 见 图 5"], border=C_AMBER, title_size=16, body_size=13)
    add_arrow(slide, 940, 750, 1020, 750)
    add_arrow(slide, 1400, 820, 1400, 880, color=C_GRAY)
    add_arrow(slide, 1400, 880, 100, 880, color=C_GRAY)
    add_arrow(slide, 100, 880, 100, 360, color=C_GRAY)
    add_label(slide, 560, 905, 500, "回到对话（闲噪 / 询问也走这条）", color=C_GRAY)
    add_footer(slide, "对话循环直到用户离开或凭证过期失效")


# ════════════════════════════════════════════════════════════
# 图5 安全闸门
# ════════════════════════════════════════════════════════════
def chart_5(slide):
    add_box(slide, 680, 230, 560, 110, "Agent 发起 request_unlock", ["credential_id + reason"], fill=C_A2, border=C_RED, accent=C_RED, title_size=16, body_size=13)
    add_arrow(slide, 960, 340, 960, 380)
    add_diamond(slide, 960, 590, 520, 400, "安全闸门五项校验",
                ["① flow = agent_confirm", "② 凭证未过期 (<300s)", "③ fusion_passed = true",
                 "④ 凭证未被消费", "⑤ 非 dry-run (NO_UNLOCK=0)"], body_size=11)
    add_arrow(slide, 1240, 590, 1420, 590, color=C_GREEN)
    add_arrow(slide, 680, 590, 500, 590, color=C_RED, dashed=True)
    add_label(slide, 1370, 580, 60, "是", color=C_GREEN, bold=True)
    add_label(slide, 585, 580, 60, "否", color=C_RED, bold=True)
    add_box(slide, 1420, 530, 420, 100, "消费凭证（幂等）", ["consumed = true · 一次性"], fill=C_A2, border=C_GREEN, accent=C_GREEN, title_size=14, body_size=12)
    add_arrow(slide, 1630, 630, 1630, 665)
    add_box(slide, 1420, 665, 420, 100, "驱动继电器 · 门开", [], fill=C_A2, border=C_GREEN, accent=C_GREEN, title_size=16)
    add_arrow(slide, 1630, 765, 1630, 800)
    add_box(slide, 1420, 800, 420, 85, "开锁完成", ["凭证立即失效 · 留痕"], fill=C_A2, border=C_GREEN, accent=C_GREEN, title_size=15, body_size=12)
    add_box(slide, 80, 530, 420, 100, "拒绝请求 X", ["返回具体拒绝原因"], border=C_RED, title_size=15, body_size=12)
    add_arrow(slide, 290, 630, 290, 665)
    add_box(slide, 80, 665, 420, 100, "Agent 告知用户", ["引导重认证 · 通知业主"], border=C_RED, title_size=15, body_size=12)
    add_arrow(slide, 290, 765, 290, 860, color=C_GRAY)
    add_arrow(slide, 290, 860, 1000, 860, color=C_GRAY)
    add_label(slide, 620, 848, 400, "回到语音对话（图4）", color=C_GRAY)
    add_box(slide, 80, 920, 860, 80, "调试保护 · 闸门只模拟开锁 · 绝不驱动继电器", [], border=C_AMBER, fill=C_CREAM, title_size=13)
    add_box(slide, 960, 920, 880, 80, "权限边界 · Agent 不评分 / 凭证过期或已消费一律拒绝", [], border=C_AMBER, fill=C_CREAM, title_size=13)
    add_footer(slide, "全流程：硬件认证 → 凭证 → 请求 → 闸门 四层防线")


def main():
    tpl = os.path.join(ROOT, "14 (1).pptx")
    prs = Presentation(tpl)
    # keep the theme, drop demo slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    # force 16:9 canvas
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    blank = prs.slide_layouts[1]  # 空白
    titles = [
        ("系统启动 · 一键拉起全线组件", "run_smart_lock.sh · 进程监督守护"),
        ("待机与触发 · 从人到认证的入口", "红外常驻轮询 · 摄像头随人开 · 自动认证条件"),
        ("五步认证 · 多模态融合判定", "人脸 45% + 活体 25% + 声纹 20% + 红外 10% · 阈值 0.78"),
        ("语音 Agent · 对话式开锁请求", "凭证就绪后开麦 · VAD → ASR → LLM → TTS 连续对话"),
        ("安全闸门 · 开锁的最终裁决", "Agent 只有请求权 · 闸门拥有一票否决权"),
    ]
    builders = [chart_1, chart_2, chart_3, chart_4, chart_5]

    for i, (fn, (title, sub)) in enumerate(zip(builders, titles), 1):
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_header(slide, i, 5, title, sub)
        fn(slide)

    out = os.path.join(ROOT, "docs", "智能门锁_流程图_5页.pptx")
    prs.save(out)
    print("saved:", out, "slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
