#!/usr/bin/env python3
"""Beautify the 3 slides of 智能门锁_.pptx with the 14(1).pptx light-blue theme.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC = "C:/Users/hoyo/Desktop/智能门锁_.pptx"
OUT = "C:/Users/hoyo/Desktop/智能门锁_美化.pptx"
BG = "C:/Users/hoyo/Desktop/lock/docs/ppt/_tpl_bg/p01_0b91503e.jpg"

SW, SH = 13.333, 7.5

C_DK = RGBColor(0x44, 0x54, 0x6A)
C_A1 = RGBColor(0xAD, 0xC7, 0xDD)
C_A2 = RGBColor(0xE1, 0xF1, 0xFE)
C_A4 = RGBColor(0x6A, 0x9A, 0xC4)
C_RED = RGBColor(0xE0, 0x62, 0x5C)
C_GREEN = RGBColor(0x4F, 0xAE, 0x7A)
C_GRAY = RGBColor(0xA6, 0xA6, 0xA6)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _no_shadow(sh):
    try:
        sh.shadow.inherit = False
    except Exception:
        pass


def _set(sh, lines, size, color, bold=False, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, font="DengXian", spacing=2):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = bold
            r.font.name = font; r.font.color.rgb = color


def add_bg(slide):
    slide.shapes.add_picture(BG, 0, 0, Inches(SW), Inches(SH))


def add_header(slide, chapter, title):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.6), Inches(0.42), Inches(3.0), Inches(0.42))
    pill.fill.solid(); pill.fill.fore_color.rgb = C_A4
    pill.line.fill.background(); _no_shadow(pill)
    _set(pill, chapter, 12, C_WHITE, bold=True)
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.94), Inches(12.1), Inches(0.5))
    _set(t, title, 24, C_DK, bold=True, align=PP_ALIGN.LEFT)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.5), Inches(1.0), Inches(0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb = C_A4
    ln.line.fill.background(); _no_shadow(ln)


def add_panel(slide, x, y, w, h, title, *, accent=C_A4):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = C_A1; card.line.width = Pt(1.5)
    _no_shadow(card)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.18), Inches(y + 0.18), Inches(0.07), Inches(0.42))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background(); _no_shadow(bar)
    tt = slide.shapes.add_textbox(Inches(x + 0.34), Inches(y + 0.14), Inches(w - 0.5), Inches(0.5))
    _set(tt, title, 17, C_DK, bold=True, align=PP_ALIGN.LEFT)
    return card


def add_item(slide, x, y, w, num, text, *, color=C_A4):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.05), Inches(0.34), Inches(0.34))
    c.fill.solid(); c.fill.fore_color.rgb = color
    c.line.fill.background(); _no_shadow(c)
    _set(c, str(num), 13, C_WHITE, bold=True)
    t = slide.shapes.add_textbox(Inches(x + 0.46), Inches(y), Inches(w - 0.46), Inches(0.5))
    _set(t, text, 14, C_DK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


# ── slides ────────────────────────────────────────────────
def p1(slide):
    add_header(slide, "立项背景与市场机会", "市场空白 与 我们的切入点")
    # left: 市场空白
    add_panel(slide, 0.6, 1.75, 6.05, 5.4, "市场空白", accent=C_RED)
    left = ["多模态多为「或逻辑」，缺自适应降级",
            "生物特征上传云端，存在隐私顾虑",
            "身高 <1.3m 儿童、>1.9m 成人覆盖不全",
            "PIR 存在误报与漏报",
            "竞品续航普遍仅 3-6 个月"]
    y = 2.55
    for i, t in enumerate(left, 1):
        add_item(slide, 0.95, y, 5.4, i, t, color=C_RED)
        y += 0.92
    # right: 切入点
    add_panel(slide, 6.8, 1.75, 5.93, 5.4, "我们的切入点", accent=C_GREEN)
    right = ["本地加权融合，阈值与权重可配置",
             "端侧比对，生物特征不出设备",
             "多模态 + 一次性凭证，判定可解释",
             "红外上升沿触发 + 状态保持，抑制抖动",
             "端侧不上云，无云存储订阅成本"]
    y = 2.55
    for i, t in enumerate(right, 1):
        add_item(slide, 7.15, y, 5.3, i, t, color=C_GREEN)
        y += 0.92


def p2(slide):
    add_header(slide, "风险边界与下一步", "项目总结")
    items = [("立项侧", "证据链完整"), ("技术侧", "板端跑通"),
             ("安全侧", "误判无法开锁"), ("工程侧", "可重复部署"),
             ("边界", "四项待验证")]
    for i, (t, b) in enumerate(items):
        add_card(slide, 0.6 + i * 2.46, 1.9, 2.3, 1.7, t, b, accent=C_A4)
    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.1), Inches(12.1), Inches(1.1))
    band.fill.solid(); band.fill.fore_color.rgb = C_A4
    band.line.fill.background(); _no_shadow(band)
    _set(band, "硬件认证负责安全 · Agent 负责交互 · 网关负责执行\n三者各司其职的组合已经验证可行", 17, C_WHITE, bold=True)


def p3(slide):
    add_header(slide, "风险边界与下一步", "当前边界与风险")
    items = [("开发板 SSH", "密钥交换阶段关闭，改动未同步"),
             ("真实继电器", "目前串口风扇模拟，需看护测试"),
             ("断网可用性", "DeepSeek 依赖网络，断网不可开门"),
             ("回声 / 延迟", "长时间打断稳定性缺正式数据"),
             ("单文件事件", "latest_event 只保留最新一条"),
             ("Bot 通知", "渠道与定时任务尚未接入"),
             ("Nano 4GB", "资源余量未知，需更小 ASR")]
    for i, (t, b) in enumerate(items):
        r, c = divmod(i, 3)
        add_card(slide, 0.6 + c * 4.1, 1.85 + r * 1.8, 3.9, 1.6, t, b,
                 accent=C_RED, title_size=15, body_size=12.5)
    f = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.4))
    _set(f, "以上为下一阶段工程工作，不影响已完成的架构验证结论", 12, C_GRAY, align=PP_ALIGN.LEFT)


def add_card(slide, x, y, w, h, title, body=None, *, accent=C_A4, title_size=16, body_size=13):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = C_WHITE
    sh.line.color.rgb = C_A1; sh.line.width = Pt(1.5)
    _no_shadow(sh)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.1), Inches(0.07), Inches(h - 0.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background(); _no_shadow(bar)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.alignment = PP_ALIGN.LEFT
    p0.space_after = Pt(3)
    for r in p0.runs:
        r.font.size = Pt(title_size); r.font.bold = True
        r.font.name = "DengXian"; r.font.color.rgb = C_DK
    if body:
        p1 = tf.add_paragraph()
        p1.text = body
        p1.alignment = PP_ALIGN.LEFT
        for r in p1.runs:
            r.font.size = Pt(body_size); r.font.name = "DengXian"; r.font.color.rgb = C_DK
    return sh


def main():
    prs = Presentation(SRC)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[1]
    for fn in (p1, p2, p3):
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        fn(slide)
    prs.save(OUT)
    print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
