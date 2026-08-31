#!/usr/bin/env python3
"""Restyle the 21 content slides of 智能门锁.pptx into the light-blue template
style. Each slide becomes structured (cards / grid / steps / timeline) instead
of a single flat block, with concise copy and larger type. Native shapes only.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
SRC = "C:/Users/hoyo/OneDrive/文档/智能门锁.pptx"
BG = os.path.join(ROOT, "docs", "ppt", "_tpl_bg", "p01_0b91503e.jpg")

SW, SH = 13.333, 7.5

# palette
C_DK = RGBColor(0x44, 0x54, 0x6A)
C_A1 = RGBColor(0xAD, 0xC7, 0xDD)
C_A2 = RGBColor(0xE1, 0xF1, 0xFE)
C_A4 = RGBColor(0x6A, 0x9A, 0xC4)
C_A5 = RGBColor(0xA3, 0xC9, 0xE9)
C_RED = RGBColor(0xE0, 0x62, 0x5C)
C_GREEN = RGBColor(0x4F, 0xAE, 0x7A)
C_GRAY = RGBColor(0xA6, 0xA6, 0xA6)
C_AMBER = RGBColor(0xE5, 0xA6, 0x4F)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CREAM = RGBColor(0xFD, 0xF5, 0xE6)


def _sz(v):
    return v if hasattr(v, "centipoints") else Pt(v)


def _no_shadow(sh):
    try:
        sh.shadow.inherit = False
    except Exception:
        pass


def _set(sh, lines, size, color, bold=False, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, font="DengXian", spacing=2):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        for r in p.runs:
            r.font.size = _sz(size)
            r.font.bold = bold
            r.font.name = font
            r.font.color.rgb = color


def add_bg(slide):
    slide.shapes.add_picture(BG, 0, 0, Inches(SW), Inches(SH))


def add_header(slide, chapter, title, num=None):
    # chapter pill
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.6), Inches(0.42), Inches(3.0), Inches(0.42))
    pill.fill.solid(); pill.fill.fore_color.rgb = C_A4
    pill.line.fill.background(); _no_shadow(pill)
    _set(pill, chapter, 12, C_WHITE, bold=True)
    # title
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.92), Inches(12.1), Inches(0.55))
    _set(t, title, 26, C_DK, bold=True, align=PP_ALIGN.LEFT)
    # accent line
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.52), Inches(1.1), Inches(0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb = C_A4
    ln.line.fill.background(); _no_shadow(ln)
    if num is not None:
        n = slide.shapes.add_textbox(Inches(12.3), Inches(6.95), Inches(0.5), Inches(0.3))
        _set(n, str(num), 9, C_GRAY, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body=None, *, fill=C_WHITE, border=C_A1,
             accent=None, title_size=16, body_size=13, title_color=C_DK, body_color=C_DK):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(1.5)
    _no_shadow(sh)
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.1), Inches(0.07), Inches(h - 0.2))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background(); _no_shadow(bar)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18 if accent else 0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.alignment = PP_ALIGN.LEFT
    p0.space_after = Pt(4)
    for r in p0.runs:
        r.font.size = _sz(title_size); r.font.bold = True
        r.font.name = "DengXian"; r.font.color.rgb = title_color
    if body:
        for ln in (body if isinstance(body, list) else [body]):
            p = tf.add_paragraph()
            p.text = ln
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(1)
            for r in p.runs:
                r.font.size = _sz(body_size); r.font.name = "DengXian"; r.font.color.rgb = body_color
    return sh


def add_stat(slide, x, y, w, h, number, label, *, color=C_A4):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = C_A1; card.line.width = Pt(1.5)
    _no_shadow(card)
    # number
    n = slide.shapes.add_textbox(Inches(x), Inches(y + 0.12), Inches(w), Inches(0.6))
    _set(n, number, 30, color, bold=True)
    # label
    l = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y + h - 0.62), Inches(w - 0.16), Inches(0.55))
    _set(l, label, 12.5, C_DK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    return card


def add_table(slide, x, y, w, headers, rows, *, col_ratios=None, font=11.5, header_font=12):
    nrows = len(rows) + 1
    ncols = len(headers)
    gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(0.4 * nrows))
    tbl = gf.table
    try:
        tbl.first_row = True
        tbl.horz_banding = False
    except Exception:
        pass
    # column widths
    if col_ratios:
        total = sum(col_ratios)
        for ci, r in enumerate(col_ratios):
            tbl.columns[ci].width = Emu(int(Inches(w) * r / total))
    else:
        for ci in range(ncols):
            tbl.columns[ci].width = Emu(int(Inches(w) / ncols))
    # header
    for ci, h in enumerate(headers):
        c = tbl.cell(0, ci)
        c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = C_A4
        c.margin_left = Inches(0.06); c.margin_right = Inches(0.06)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = _sz(header_font); r.font.bold = True
                r.font.name = "DengXian"; r.font.color.rgb = C_WHITE
    # body
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.text = val
            c.margin_left = Inches(0.06); c.margin_right = Inches(0.06)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = C_A2
            else:
                c.fill.solid(); c.fill.fore_color.rgb = C_WHITE
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = _sz(font); r.font.name = "DengXian"; r.font.color.rgb = C_DK
    return tbl


def add_step_h(slide, num, x, y, w, h, text, *, color=C_A4):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = C_WHITE
    sh.line.color.rgb = C_A1; sh.line.width = Pt(1.5)
    _no_shadow(sh)
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.14), Inches(y + 0.14), Inches(0.5), Inches(0.5))
    c.fill.solid(); c.fill.fore_color.rgb = color
    c.line.fill.background(); _no_shadow(c)
    _set(c, str(num), 18, C_WHITE, bold=True)
    t = slide.shapes.add_textbox(Inches(x + 0.72), Inches(y), Inches(w - 0.82), Inches(h))
    _set(t, text, 13, C_DK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    return sh


def add_arrow(slide, x1, y1, x2, y2, color=C_A4, width=2.5):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "none", "w": "med", "len": "med"}))
    return conn


def add_timeline(slide, x, y, w, milestones, *, label_size=13, desc_size=10.5):
    n = len(milestones)
    line_y = y + 0.9
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.3), Inches(line_y), Inches(w - 0.6), Inches(0.035))
    ln.fill.solid(); ln.fill.fore_color.rgb = C_A1
    ln.line.fill.background(); _no_shadow(ln)
    step = (w - 0.6) / n
    for i, ms in enumerate(milestones):
        cx = x + 0.3 + step * i + step / 2
        label = ms.get("label", ""); desc = ms.get("desc", "")
        active = ms.get("active", False)
        col = C_A4 if active else C_A1
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.11), Inches(line_y - 0.095), Inches(0.22), Inches(0.22))
        dot.fill.solid(); dot.fill.fore_color.rgb = col
        dot.line.fill.background(); _no_shadow(dot)
        lb = slide.shapes.add_textbox(Inches(cx - step / 2 + 0.04), Inches(y), Inches(step - 0.08), Inches(0.55))
        _set(lb, label, label_size, C_DK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        if desc:
            db = slide.shapes.add_textbox(Inches(cx - step / 2 + 0.04), Inches(line_y + 0.25), Inches(step - 0.08), Inches(0.6))
            _set(db, desc, desc_size, C_DK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


# ── slide builders ─────────────────────────────────────────
def s01(slide):
    add_header(slide, "立项背景与市场机会", "市场结构：渗透率、品类与集中度", 6)
    stats = [("<30%", "城镇家庭渗透率\n一线 45% · 三四线 22%"),
             (">60%", "日韩渗透率\n中国仍处普及中段"),
             ("43.2%", "人脸识别\n传统电商销量占比"),
             ("11.2%", "AI 智能锁\n2025 销量占比"),
             ("28.3%", "TOP3 品牌\n小米·德施曼·凯迪仕"),
             ("450–700 亿", "TAM 市场空间\n每年")]
    x0, y0, gw, gh, gap = 0.6, 1.85, 3.9, 1.6, 0.21
    for i, (num, lab) in enumerate(stats):
        r, c = divmod(i, 3)
        add_stat(slide, x0 + c * (gw + gap), y0 + r * (gh + 0.25), gw, gh, num, lab)
    # footnote
    f = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.35))
    _set(f, "数据口径统一采用 01_ 市场机会分析报告", 11, C_GRAY, align=PP_ALIGN.LEFT)


def s02(slide):
    add_header(slide, "立项背景与市场机会", "竞品格局对比", 9)
    add_table(slide, 1.2, 1.9, 10.9, ["竞品", "主要配置", "价格带"],
              [["德施曼 R8", "双目红外 3D 人脸 + 指静脉 + 大屏猫眼", "3000-4000 元"],
               ["凯迪仕 K70", "3D 人脸 + 掌静脉 + 本地 AI 语音管家", "3500-5000 元"],
               ["小米 M30 Pro", "AI 3D 结构光 + 广角猫眼 + 米家生态", "2000-3000 元"],
               ["华为 Pro", "3D ToF + 分布式可视猫眼 + 鸿蒙", "2500-3500 元"]],
              col_ratios=[1.5, 4.5, 1.8], font=13, header_font=13)


def s03(slide):
    add_header(slide, "关键技术预研与选型", "技术选型：Demo 验证 vs 量产方案", 13)
    add_table(slide, 1.0, 1.9, 11.3, ["能力域", "Demo 验证路线", "量产方案路线"],
              [["深度估计", "Depth Anything / MiDaS", "3D 结构光优先，双目 IR 备选"],
               ["人脸识别", "InsightFace / ArcFace", "3D 结构光 + 端侧比对"],
               ["活体检测", "MediaPipe 动作活体", "3D 活体 + 分级唤醒"],
               ["声纹识别", "SpeechBrain ECAPA", "声纹作为第二因子"],
               ["主认证兜底", "规则引擎 + 人脸声纹融合", "半导体指纹 + 密码兜底"],
               ["多模态融合", "规则引擎加权融合", "规则化多模态，策略从简"],
               ["执行与推理", "ONNX Runtime / TensorRT", "全自动锁体 + 端侧 NPU"]],
              col_ratios=[1.4, 3.2, 3.4], font=12, header_font=13)


def s04(slide):
    add_header(slide, "关键技术预研与选型", "弃用技术路线及理由", 14)
    items = [("常电视觉", "功耗高、续航差"),
             ("云端比对", "隐私与断网风险"),
             ("光学指纹", "干湿手识别弱"),
             ("复杂 AI 融合", "可解释性差 · 留 V2.0"),
             ("UWB / 声纹开锁", "可行性低 · 仅第二因子"),
             ("指静脉", "识别慢、体积大")]
    x0, y0, cw, ch, gx, gy = 0.6, 1.85, 3.9, 1.5, 0.21, 0.25
    for i, (t, b) in enumerate(items):
        r, c = divmod(i, 3)
        add_card(slide, x0 + c * (cw + gx), y0 + r * (ch + gy), cw, ch, t, b,
                 accent=C_A4, title_size=17, body_size=13.5)


def s05(slide):
    add_header(slide, "关键技术预研与选型", "关键技术项判定结论", 15)
    add_table(slide, 1.0, 1.9, 11.3, ["技术项", "判定", "说明"],
              [["人脸识别", "Pass", "Demo 用 2D，量产转 3D 结构光"],
               ["指纹识别", "Pass", "量产采用半导体指纹"],
               ["密码兜底", "Pass", "保留为最终降级通道"],
               ["声纹识别", "Conditional", "仅作第二因子，不单独放行"],
               ["多模态融合", "Pass", "框架先行，策略从简"],
               ["低功耗", "Pass（设计验证）", "分级唤醒，未做真机续航验证"],
               ["安全", "Pass", "端侧比对，特征不出设备"],
               ["锁体", "Pass（方案级）", "全自动锁体仅方案分析"],
               ["遮挡场景", "Conditional", "准确率 72%→95.3%，仍需实测"]],
              col_ratios=[1.4, 2.0, 4.6], font=12, header_font=13)


def s06(slide):
    add_header(slide, "系统架构与实现", "系统四层职责边界", 19)
    items = [("Jetson 硬件认证", "独立完成多模态融合判定，只产生凭证"),
             ("FastGPT + DeepSeek", "只做语言理解与工具编排"),
             ("Jetson 工具网关", "只验证凭证与干运行，驱动执行器"),
             ("Pipecat 语音流水线", "只做音频、打断与轮流说话"),
             ("异常事件服务", "只记录失败，不干预判定")]
    # 2+3 grid
    coords = [(0.6, 1.9, 7.3, 1.5), (8.1, 1.9, 4.6, 1.5),
              (0.6, 3.65, 3.9, 1.6), (4.7, 3.65, 3.9, 1.6), (8.8, 3.65, 3.9, 1.6)]
    for (t, b), (x, y, w, h) in zip(items, coords):
        add_card(slide, x, y, w, h, t, b, accent=C_A4, title_size=16, body_size=12.5)


def s07(slide):
    add_header(slide, "系统架构与实现", "硬件多模态认证流程", 20)
    steps = ["红外检测\n人体靠近", "人脸识别\nInsightFace", "活体检测\n眨眼 + 转头",
             "声纹识别\n文本无关", "融合判定\n加权求和"]
    x0, sw_, sh_, gap = 0.6, 2.25, 2.5, 0.24
    for i, txt in enumerate(steps, 1):
        x = x0 + (i - 1) * (sw_ + gap)
        add_step_h(slide, i, x, 2.0, sw_, sh_, txt)
        if i < 5:
            add_arrow(slide, x + sw_, 3.25, x + sw_ + gap, 3.25)
    f = slide.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.4))
    _set(f, "权重：人脸 0.45 · 活体 0.25 · 声纹 0.20 · 红外 0.10　阈值 0.78　（详见下页）",
         13, C_A4, align=PP_ALIGN.LEFT)


def s08(slide):
    add_header(slide, "系统架构与实现", "多模态融合权重与阈值", 21)
    add_table(slide, 1.2, 1.9, 10.9, ["模块", "融合权重", "单项阈值与关键参数"],
              [["人脸（InsightFace）", "0.45", "min_score 0.65"],
               ["活体（MediaPipe）", "0.25", "min_score 0.60 · 眨眼 + 转头"],
               ["声纹（ECAPA）", "0.20", "min_score 0.70 · 最短语音 2s"],
               ["红外（串口）", "0.10", "上升沿触发 · 保持 30s"],
               ["融合阈值", "0.78", "require_all = false · 加权总分"]],
              col_ratios=[2.4, 1.4, 4.2], font=13, header_font=13)
    # weight bar
    add_header  # no-op guard
    bar = slide.shapes.add_textbox(Inches(1.2), Inches(5.9), Inches(10.9), Inches(0.4))
    _set(bar, "各模块权重按 config.yaml 可调，Agent 不参与评分", 11.5, C_GRAY, align=PP_ALIGN.LEFT)


def s09(slide):
    add_header(slide, "系统架构与实现", "低延迟语音 Agent 流水线", 24)
    steps = ["麦克风\n16 kHz 采集", "VAD\n起止 / 打断", "ASR\n本地识别",
             "LLM\n流式回复", "TTS\n本地合成", "播放\naplay 输出"]
    x0, sw_, sh_, gap = 0.45, 1.95, 2.35, 0.14
    for i, txt in enumerate(steps, 1):
        x = x0 + (i - 1) * (sw_ + gap)
        add_step_h(slide, i, x, 2.0, sw_, sh_, txt, color=C_A4 if i not in (4, 5) else C_A5)
        if i < 6:
            add_arrow(slide, x + sw_, 3.18, x + sw_ + gap, 3.18)
    f = slide.shapes.add_textbox(Inches(0.45), Inches(5.2), Inches(12.4), Inches(0.4))
    _set(f, "ASR / TTS 本地执行，DeepSeek 依赖网络 → 非完全离线", 12, C_GRAY, align=PP_ALIGN.LEFT)


def s10(slide):
    add_header(slide, "系统架构与实现", "Agent 受控工具（仅两个）", 25)
    rows = [["current_auth_context", "GET /tools/current_auth_context", "读取最新硬件认证凭证", "开锁前必须先调用"],
            ["request_unlock", "POST /tools/request_unlock", "请求开锁，闸门裁决", "四条件同时满足才调用"]]
    add_table(slide, 1.0, 1.9, 11.3, ["工具", "方法 / 路径", "用途", "调用约束"], rows,
              col_ratios=[2.2, 3.4, 2.4, 3.0], font=13, header_font=13)
    f = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.4))
    _set(f, "不注册裸 unlock，Agent 不能评分、否决或伪造硬件认证", 12, C_RED, align=PP_ALIGN.LEFT)


def s11(slide):
    add_header(slide, "安全设计与异常闭环", "四层安全约束", 28)
    items = [("第一层 · 融合层", "只有本地 FusionEngine 产生认证结果"),
             ("第二层 · 凭证层", "短时一次性，过期重复消费一律拒绝"),
             ("第三层 · Agent 层", "仅两个受控工具，不暴露裸 unlock"),
             ("第四层 · 执行层", "Bearer Token + 干运行 + 本地执行器")]
    for i, (t, b) in enumerate(items):
        r, c = divmod(i, 2)
        add_card(slide, 0.6 + c * 6.2, 1.9 + r * 1.85, 6.0, 1.6, t, b,
                 accent=C_A4, title_size=17, body_size=13.5)
    # conclusion band
    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.75), Inches(12.1), Inches(0.95))
    band.fill.solid(); band.fill.fore_color.rgb = C_A2
    band.line.color.rgb = C_A4; band.line.width = Pt(1.5); _no_shadow(band)
    _set(band, "即使模型误解或提示词失效，本地闸门仍能拒绝无凭证请求", 15, C_DK, bold=True)


def s12(slide):
    add_header(slide, "安全设计与异常闭环", "安全验证：两个真实场景", 29)
    add_card(slide, 0.6, 1.9, 6.0, 3.2, "有效凭证 + 「请开门」",
             ["current_auth_context → request_unlock", "结果：干运行阻止真实动作"],
             accent=C_GREEN, title_size=18, body_size=14)
    add_card(slide, 6.9, 1.9, 5.8, 3.2, "无效凭证 + 「请开门」",
             ["仅调用 current_auth_context", "结果：不调用开锁工具"],
             accent=C_RED, title_size=18, body_size=14)
    f = slide.shapes.add_textbox(Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.5))
    _set(f, "结论：模型无法绕过本地安全闸门", 14, C_DK, bold=True, align=PP_ALIGN.LEFT)


def s13(slide):
    add_header(slide, "安全设计与异常闭环", "端口与服务分工", 31)
    add_table(slide, 1.0, 1.9, 11.3, ["服务", "端口", "部署位置", "配置入口"],
              [["FastGPT", "3300", "Windows / WSL", "FASTGPT_PORT / API_BASE"],
               ["门锁工具网关", "8787", "Jetson", "LOCK_TOOL_GATEWAY_PORT"],
               ["异常事件服务", "8790", "Jetson 本机", "LOCK_EVENT_SERVICE_PORT"]],
              col_ratios=[1.6, 1.0, 2.2, 3.6], font=13.5, header_font=13)


def s14(slide):
    add_header(slide, "安全设计与异常闭环", "进程治理与一键运维", 32)
    items = [("一键脚本", "start / stop / restart"),
             ("统一监督", "GUI · Agent · 网关 · 事件"),
             ("异常退出", "自动拉起"),
             ("停止清理", "PID · 临时文件 · 端口"),
             ("幂等启动", "清旧凭证"),
             ("SSH 隧道", "Windows 侧管理")]
    for i, (t, b) in enumerate(items):
        r, c = divmod(i, 3)
        add_card(slide, 0.6 + c * 4.1, 1.9 + r * 2.15, 3.9, 1.9, t, b,
                 accent=C_A4, title_size=16, body_size=13.5)


def s15(slide):
    add_header(slide, "验证结果与工程问题", "测试与验证结果清单", 34)
    rows = [["原硬件正向认证流程", "开发板已跑通", C_GREEN],
            ["FastGPT Web 可达性", "localhost:3300 → 200", C_GREEN],
            ["Agent 有效/无效凭证工具调用", "两场景均通过", C_GREEN],
            ["工具网关生命周期", "start/stop/restart/crash 通过", C_GREEN],
            ["异常事件服务生命周期", "start/stop/restart/crash 通过", C_GREEN],
            ["端口与父进程环境清理", "通过", C_GREEN],
            ["Pipecat 单元链路", "通过", C_GREEN],
            ["Windows/WSL 文件语音链路", "通过", C_GREEN],
            ["开发板真麦克风音箱对话", "历史版本通过", C_GREEN],
            ["最新异常触发 + aplay 输出", "待板端复测", C_AMBER]]
    # two-column list of status chips
    add_table(slide, 1.0, 1.9, 11.3, ["测试项", "结果"], [[r[0], r[1]] for r in rows],
              col_ratios=[6.0, 3.0], font=12.5, header_font=13)


def s16(slide):
    add_header(slide, "验证结果与工程问题", "关键工程问题与解决方案", 36)
    add_table(slide, 0.9, 1.9, 11.5, ["问题", "解决方案"],
              [["红外持续 true 重复播报", "上升沿触发 + 状态保持 + 离开重置"],
               ["模型首次使用延迟高", "GUI 与 FunASR 启动时预加载"],
               ["声纹与 Agent 抢麦克风", "Agent 等待凭证，认证后开麦"],
               ["Agent 可能绕过认证", "一次性凭证 + 本地网关"],
               ["FastGPT 无法访问 Jetson", "SSH 本地 / 反向隧道"],
               ["异常退出污染端口", "PID 监督 + 信号清理 + 恢复测试"],
               ["异常情况无记录", "独立事件服务 + 单文件原子写入"],
               ["音频路径不统一", "输入 PortAudio · 输出 aplay"]],
              col_ratios=[1.0, 1.0], font=12.5, header_font=13)


def s17(slide):
    add_header(slide, "商业规划与产品路标", "成本、定价与营收测算", 38)
    add_table(slide, 0.9, 1.9, 11.5, ["项目", "Demo 阶段", "量产测算（教学示例）"],
              [["硬件 / 耗材", "0-2000 元（不含已有设备）", "边际成本 1000-1600 元/套"],
               ["人力与周期", "5-8 人 · 3-4 个月", "—"],
               ["毛利率", "—", "30%-45%"],
               ["建议定价", "—", "2599 / 2999 / 3499 元"],
               ["首年销量营收", "—", "1-3 万套 · 0.3-0.9 亿元"],
               ["研发营销投入", "可选用免费云额度", "2000-5000 万元"],
               ["投资回收期", "—", "2-4 年"]],
              col_ratios=[1.5, 3.0, 3.5], font=12.5, header_font=13)


def s18(slide):
    add_header(slide, "商业规划与产品路标", "版本路标", 40)
    add_timeline(slide, 0.7, 2.2, 11.9, [
        {"label": "V1.0", "desc": "Demo 与量产基础能力", "active": True},
        {"label": "V2.0", "desc": "多模态自适应\n遮挡优化"},
        {"label": "V2.0+", "desc": "智能猫眼\nMatter / HomeKit"},
        {"label": "V3.0", "desc": "声纹独立开锁\nAI 主动安防"},
        {"label": "V3.0+", "desc": "指静脉 / UWB\n统一身份平台"}],
        label_size=15, desc_size=11)


def s19(slide):
    add_header(slide, "商业规划与产品路标", "IPD 流程节点与交付物", 41)
    add_timeline(slide, 0.7, 2.2, 11.9, [
        {"label": "CDCP", "desc": "Charter 审批\nPDT 组建", "active": True},
        {"label": "TR1", "desc": "需求说明\n产品概念"},
        {"label": "TR2 / TR3", "desc": "架构选型\n详细设计"},
        {"label": "TR4 / TR5", "desc": "模块测试\n整机联调"},
        {"label": "TR6 / ADCP", "desc": "归档与结项\n演示（2026-12）"}],
        label_size=14, desc_size=10.5)


def s20(slide):
    add_header(slide, "风险边界与下一步", "当前边界与风险", 43)
    items = ["开发板 SSH 待恢复", "真实继电器待测", "断网对话不可用",
             "回声/延迟缺数据", "单文件事件队列", "Bot 通知未接入",
             "Nano 4GB 未压测"]
    for i, txt in enumerate(items):
        r, c = divmod(i, 3)
        add_card(slide, 0.6 + c * 4.1, 1.9 + r * 1.75, 3.9, 1.55, txt, None,
                 accent=C_RED, title_size=15)
    f = slide.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.4))
    _set(f, "以上为下一阶段工程工作，不影响已完成的架构验证结论", 12, C_GRAY, align=PP_ALIGN.LEFT)


def s21(slide):
    add_header(slide, "风险边界与下一步", "项目总结", 45)
    items = [("立项侧", "证据链完整"), ("技术侧", "板端跑通"),
             ("安全侧", "误判无法开锁"), ("工程侧", "可重复部署"),
             ("边界", "四项待验证")]
    for i, (t, b) in enumerate(items):
        add_card(slide, 0.6 + i * 2.46, 1.9, 2.3, 1.7, t, b,
                 accent=C_A4, title_size=17, body_size=13.5)
    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.1), Inches(12.1), Inches(1.1))
    band.fill.solid(); band.fill.fore_color.rgb = C_A4
    band.line.fill.background(); _no_shadow(band)
    _set(band, "硬件认证负责安全 · Agent 负责交互 · 网关负责执行\n三者各司其职的组合已经验证可行", 17, C_WHITE, bold=True)


def main():
    prs = Presentation(SRC)
    # keep theme; wipe demo slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[1]

    builders = [s01, s02, s03, s04, s05, s06, s07, s08, s09, s10, s11,
                s12, s13, s14, s15, s16, s17, s18, s19, s20, s21]
    for fn in builders:
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        fn(slide)

    out = "C:/Users/hoyo/OneDrive/文档/智能门锁_美化版.pptx"
    prs.save(out)
    print("saved:", out, "slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    from pptx.util import Emu
    main()
