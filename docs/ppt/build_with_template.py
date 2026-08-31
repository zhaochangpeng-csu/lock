#!/usr/bin/env python3
"""
Build the detailed project deck on top of the light-blue business template
"14 (1).pptx" (10 x 5.625 in). Keeps the template's slide size, theme colors
and background gradient, while replacing all 23 demo slides with the 46-page
report from slides.json.
"""
import json
import os
import sys
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData


# Theme colors extracted from 14 (1).pptx
COLORS = {
    "dk1": RGBColor(0x00, 0x00, 0x00),
    "lt1": RGBColor(0xFF, 0xFF, 0xFF),
    "dk2": RGBColor(0x44, 0x54, 0x6A),
    "lt2": RGBColor(0xE6, 0xE4, 0xE4),
    "accent1": RGBColor(0xAD, 0xC7, 0xDD),
    "accent2": RGBColor(0xE1, 0xF1, 0xFE),
    "accent3": RGBColor(0xAF, 0xC8, 0xDE),
    "accent4": RGBColor(0x6A, 0x9A, 0xC4),
    "accent5": RGBColor(0xA3, 0xC9, 0xE9),
    "accent6": RGBColor(0x6E, 0xAC, 0x46),
}

SLIDE_W = 10.0
SLIDE_H = 5.625

BG_IMAGE = os.path.join(os.path.dirname(__file__), "_tpl_bg", "p01_0b91503e.jpg")


def _c(name):
    return COLORS.get(name, COLORS["dk2"])


def _add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def _set_font(run, size, bold=False, color=None, name="DengXian"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    # For Chinese rendering, also set east Asian typeface explicitly.
    try:
        rPr = run.font._element
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        ea = rPr.find(f"{ns}ea")
        if ea is None:
            from xml.etree.ElementTree import SubElement
            SubElement(rPr, f"{ns}ea")
            ea = rPr.find(f"{ns}ea")
        ea.set("typeface", name)
    except Exception:
        pass
    if color:
        run.font.color.rgb = color


def _add_para(tf, text, size=18, bold=False, color=None, align=None, spacing=0, level=0, font_name="DengXian"):
    if tf.paragraphs and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.level = level
    if align:
        p.alignment = align
    if spacing:
        p.space_after = Pt(spacing)
    # Apply font to all runs (sometimes text splits into multiple runs)
    for r in p.runs:
        _set_font(r, size, bold=bold, color=color, name=font_name)
    # If no runs were created (e.g. empty), set the paragraph font.
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return p


def _add_filled_shape(slide, shape_type, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_gradient_bg(slide):
    slide.shapes.add_picture(BG_IMAGE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))


def _add_content_card(slide):
    # White rounded card for content slides (ends at y=5.25 to leave a footer band).
    card = _add_filled_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                             0.55, 0.85, 8.9, 4.40, _c("lt1"))
    card.line.color.rgb = _c("accent1")
    card.line.width = Pt(0.75)
    return card


def _add_page_footer(slide, page, footer_text="智能门锁 AI 原型 · IPD 立项与技术汇报"):
    # Page number
    box = _add_textbox(slide, 0.50, 5.30, 0.60, 0.25)
    _add_para(box.text_frame, str(page), size=7, color=_c("dk2"), align=PP_ALIGN.LEFT)
    # Footer text
    if footer_text:
        box = _add_textbox(slide, 5.50, 5.30, 3.95, 0.25)
        _add_para(box.text_frame, footer_text, size=7, color=_c("dk2"), align=PP_ALIGN.RIGHT)


class DeckBuilder:
    def __init__(self, template_path):
        self.prs = Presentation(template_path)
        # Ensure the background asset exists.
        if not os.path.exists(BG_IMAGE):
            raise FileNotFoundError(f"Missing template background: {BG_IMAGE}")
        # Use layout "空白" (index 1) for new slides.
        self.blank_layout = self.prs.slide_layouts[1]
        self._page = 0
        # Remove all existing demo slides.
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

    def _new_slide(self):
        self._page += 1
        return self.prs.slides.add_slide(self.blank_layout)

    def build_cover(self, s):
        slide = self._new_slide()
        _add_gradient_bg(slide)
        title = s.get("title", "")
        subtitle = s.get("subtitle", "")
        author = s.get("author", "")
        date = s.get("date", "")
        # Title
        tb = _add_textbox(slide, 1.25, 1.85, 7.5, 1.0)
        _add_para(tb.text_frame, title, size=36, bold=True, color=_c("dk2"),
                  align=PP_ALIGN.CENTER, spacing=6)
        # Subtitle
        if subtitle:
            tb = _add_textbox(slide, 1.25, 3.10, 7.5, 0.8)
            for line in subtitle.split("\n"):
                _add_para(tb.text_frame, line, size=15, color=_c("dk2"),
                          align=PP_ALIGN.CENTER, spacing=4)
        # Meta line
        meta = "  |  ".join(filter(None, [author, date]))
        if meta:
            tb = _add_textbox(slide, 1.25, 4.50, 7.5, 0.35)
            _add_para(tb.text_frame, meta, size=11, color=_c("accent4"),
                      align=PP_ALIGN.CENTER)
        _add_page_footer(slide, self._page, "")
        return slide

    def build_section(self, s):
        slide = self._new_slide()
        _add_gradient_bg(slide)
        # Accent bar
        bar = _add_filled_shape(slide, MSO_SHAPE.RECTANGLE, 1.20, 2.05, 0.05, 1.35, _c("accent4"))
        bar.line.fill.background()
        # Big chapter number (use the slide number as the section index)
        num = str(self._page)
        tb = _add_textbox(slide, 6.50, 1.05, 2.8, 3.2)
        _add_para(tb.text_frame, num, size=115, bold=True, color=_c("accent4"),
                  align=PP_ALIGN.RIGHT, font_name="Agency FB")
        # Title
        tb = _add_textbox(slide, 1.45, 2.25, 6.8, 0.7)
        _add_para(tb.text_frame, s.get("title", ""), size=26, bold=True, color=_c("dk2"),
                  align=PP_ALIGN.LEFT, spacing=4)
        # Subtitle
        subtitle = s.get("subtitle", "")
        if subtitle:
            tb = _add_textbox(slide, 1.45, 3.00, 6.8, 0.6)
            _add_para(tb.text_frame, subtitle, size=13, color=_c("dk2"), align=PP_ALIGN.LEFT)
        _add_page_footer(slide, self._page)
        return slide

    def build_agenda(self, s):
        slide = self._new_slide()
        _add_gradient_bg(slide)
        # "目录" title
        tb = _add_textbox(slide, 0.85, 1.95, 3.0, 1.0)
        _add_para(tb.text_frame, "目录", size=48, bold=True, color=_c("dk2"), align=PP_ALIGN.LEFT)
        tb = _add_textbox(slide, 0.85, 3.05, 3.0, 0.3)
        _add_para(tb.text_frame, "CONTENTS", size=11, color=_c("accent4"), align=PP_ALIGN.LEFT)
        # Topics as rounded labels on the right
        topics = s.get("topics", [])
        y0 = 1.05
        for i, t in enumerate(topics, 1):
            y = y0 + (i - 1) * 0.62
            # Number circle
            circle = _add_filled_shape(slide, MSO_SHAPE.OVAL, 4.05, y, 0.34, 0.34, _c("accent4"))
            ct = _add_textbox(slide, 4.05, y + 0.02, 0.34, 0.30)
            _add_para(ct.text_frame, f"{i:02d}", size=10, bold=True, color=_c("lt1"),
                      align=PP_ALIGN.CENTER)
            # Text
            tt = _add_textbox(slide, 4.55, y + 0.03, 4.7, 0.30)
            _add_para(tt.text_frame, t, size=15, color=_c("dk2"), align=PP_ALIGN.LEFT)
        _add_page_footer(slide, self._page)
        return slide

    def _card_title(self, slide, title):
        _add_content_card(slide)
        tb = _add_textbox(slide, 0.85, 0.95, 8.3, 0.55)
        _add_para(tb.text_frame, title, size=22, bold=True, color=_c("accent4"), align=PP_ALIGN.LEFT)
        # Underline accent
        line = _add_filled_shape(slide, MSO_SHAPE.RECTANGLE, 0.85, 1.45, 8.3, 0.02, _c("accent1"))
        line.line.fill.background()

    def build_content(self, s):
        slide = self._new_slide()
        _add_gradient_bg(slide)
        self._card_title(slide, s.get("title", ""))
        items = s.get("items", [])
        tb = _add_textbox(slide, 0.85, 1.60, 8.3, 3.45)
        for item in items:
            if isinstance(item, dict):
                _add_para(tb.text_frame, item.get("t", ""), size=15, bold=True,
                          color=_c("dk2"), spacing=2)
                for sub in item.get("s", []):
                    _add_para(tb.text_frame, sub, size=12, color=_c("dk2"),
                              spacing=0, level=1)
            else:
                _add_para(tb.text_frame, str(item), size=13, color=_c("dk2"), spacing=6)
        _add_page_footer(slide, self._page)
        return slide

    def build_two_column(self, s):
        slide = self._new_slide()
        _add_gradient_bg(slide)
        self._card_title(slide, s.get("title", ""))
        # Divider
        _add_filled_shape(slide, MSO_SHAPE.RECTANGLE, 4.975, 1.80, 0.02, 3.0, _c("accent1"))
        left = s.get("left", [])
        right = s.get("right", [])
        for col, x, items in [("L", 0.85, left), ("R", 5.15, right)]:
            tb = _add_textbox(slide, x, 1.75, 3.9, 3.25)
            for item in items:
                _add_para(tb.text_frame, str(item), size=12.5, color=_c("dk2"), spacing=7)
        _add_page_footer(slide, self._page)
        return slide

    def build_table(self, s):
        slide = self._new_slide()
        _add_gradient_bg(slide)
        self._card_title(slide, s.get("title", ""))
        headers = s.get("headers", [])
        rows = s.get("rows", [])
        nrows = len(rows) + 1
        ncols = len(headers)
        if ncols == 0 or nrows == 0:
            _add_page_footer(slide, self._page)
            return slide
        left, top = 0.55, 1.70
        width = 8.9
        height = 3.45
        tbl = slide.shapes.add_table(nrows, ncols, Inches(left), Inches(top),
                                     Inches(width), Inches(height)).table
        # Header
        for i, h in enumerate(headers):
            c = tbl.cell(0, i)
            c.text = str(h)
            c.fill.solid()
            c.fill.fore_color.rgb = _c("accent4")
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = _c("lt1")
                p.alignment = PP_ALIGN.CENTER
        # Rows
        for ri, row in enumerate(rows, 1):
            for ci, val in enumerate(row):
                c = tbl.cell(ri, ci)
                c.text = str(val)
                if ri % 2 == 0:
                    c.fill.solid()
                    c.fill.fore_color.rgb = _c("accent2")
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.font.color.rgb = _c("dk2")
                    p.alignment = PP_ALIGN.CENTER
        _add_page_footer(slide, self._page)
        return slide

    def build_image_slide(self, s):
        slide = self._new_slide()
        _add_gradient_bg(slide)
        self._card_title(slide, s.get("title", ""))
        image_path = s.get("image_path", "")
        if not os.path.exists(image_path):
            # Try relative to project root.
            image_path = os.path.join("C:/Users/hoyo/Desktop/lock", image_path)
        if os.path.exists(image_path):
            # Fit-contain inside the card area.
            box_x, box_y, box_w, box_h = 0.70, 1.75, 8.6, 3.30
            try:
                from PIL import Image
                with Image.open(image_path) as im:
                    iw, ih = im.size
                scale = min(box_w / iw, box_h / ih)
                dw, dh = iw * scale, ih * scale
                px = box_x + (box_w - dw) / 2
                py = box_y + (box_h - dh) / 2
            except Exception:
                px, py, dw, dh = box_x, box_y, box_w, box_h
            slide.shapes.add_picture(image_path, Inches(px), Inches(py), Inches(dw), Inches(dh))
        caption = s.get("caption", "")
        if caption:
            tb = _add_textbox(slide, 0.70, 5.05, 8.6, 0.20)
            _add_para(tb.text_frame, caption, size=9, color=_c("dk2"), align=PP_ALIGN.CENTER)
        _add_page_footer(slide, self._page)
        return slide

    def build_summary(self, s):
        slide = self._new_slide()
        _add_gradient_bg(slide)
        self._card_title(slide, s.get("title", ""))
        points = s.get("points", [])
        tb = _add_textbox(slide, 0.95, 1.75, 8.2, 2.35)
        for pt in points:
            txt = pt if isinstance(pt, str) else pt.get("t", str(pt))
            _add_para(tb.text_frame, txt, size=13, color=_c("dk2"), spacing=10)
        conclusion = s.get("conclusion", "")
        if conclusion:
            box = _add_filled_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                                    0.85, 4.15, 8.3, 0.75, _c("accent4"))
            tf = box.text_frame
            tf.word_wrap = True
            _add_para(tf, conclusion, size=13, bold=True, color=_c("lt1"),
                      align=PP_ALIGN.CENTER)
        _add_page_footer(slide, self._page)
        return slide

    def build_timeline(self, s):
        slide = self._new_slide()
        _add_gradient_bg(slide)
        self._card_title(slide, s.get("title", ""))
        milestones = s.get("milestones", [])
        n = len(milestones)
        if n == 0:
            _add_page_footer(slide, self._page)
            return slide
        spacing = 8.4 / max(n, 1)
        line_y = 3.4
        # Horizontal line
        _add_filled_shape(slide, MSO_SHAPE.RECTANGLE, 0.80, line_y, 8.4, 0.03, _c("accent1"))
        for i, ms in enumerate(milestones):
            cx = 0.80 + spacing * i + spacing / 2
            label = ms.get("label", "") if isinstance(ms, dict) else str(ms)
            desc = ms.get("desc", "") if isinstance(ms, dict) else ""
            active = ms.get("active", False) if isinstance(ms, dict) else False
            dot_color = _c("accent4") if active else _c("accent1")
            _add_filled_shape(slide, MSO_SHAPE.OVAL, cx - 0.12, line_y - 0.105, 0.24, 0.24, dot_color)
            lb = _add_textbox(slide, cx - spacing * 0.45, line_y - 0.85, spacing * 0.9, 0.45)
            _add_para(lb.text_frame, label, size=12, bold=True, color=_c("dk2"),
                      align=PP_ALIGN.CENTER, spacing=2)
            if desc:
                db = _add_textbox(slide, cx - spacing * 0.45, line_y + 0.35, spacing * 0.9, 0.50)
                _add_para(db.text_frame, desc, size=9, color=_c("dk2"), align=PP_ALIGN.CENTER)
        _add_page_footer(slide, self._page)
        return slide

    def build_chart(self, s):
        slide = self._new_slide()
        _add_gradient_bg(slide)
        self._card_title(slide, s.get("title", ""))
        cd = CategoryChartData()
        cd.categories = s.get("categories", [])
        for series in s.get("series", []):
            cd.add_series(series.get("name", ""), series.get("values", []))
        ct_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }
        ct = ct_map.get(s.get("chart_type", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
        chart_shape = slide.shapes.add_chart(ct, Inches(0.80), Inches(1.70),
                                              Inches(8.4), Inches(3.45), cd)
        chart = chart_shape.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        if hasattr(chart, 'font'):
            chart.font.size = Pt(9)
        if s.get("chart_type") == "pie":
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.font.size = Pt(9)
        _add_page_footer(slide, self._page)
        return slide

    def build_quote(self, s):
        slide = self._new_slide()
        _add_gradient_bg(slide)
        self._card_title(slide, s.get("title", ""))
        box = _add_filled_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                                1.0, 1.85, 8.0, 2.6, _c("accent2"))
        qm = _add_textbox(slide, 1.0, 1.75, 0.8, 0.6)
        _add_para(qm.text_frame, "❝", size=36, bold=True, color=_c("accent4"),
                  align=PP_ALIGN.CENTER)
        tb = _add_textbox(slide, 1.7, 2.1, 6.9, 2.0)
        _add_para(tb.text_frame, s.get("quote", ""), size=18, color=_c("dk2"),
                  align=PP_ALIGN.LEFT, spacing=8)
        attribution = s.get("attribution", "")
        if attribution:
            ab = _add_textbox(slide, 1.7, 4.1, 6.9, 0.3)
            _add_para(ab.text_frame, f"— {attribution}", size=12, color=_c("dk2"),
                      align=PP_ALIGN.RIGHT)
        _add_page_footer(slide, self._page)
        return slide

    def build_contact(self, s):
        slide = self._new_slide()
        _add_gradient_bg(slide)
        title = s.get("title", "")
        info = s.get("info", "")
        if isinstance(info, list):
            info = "\n".join(info)
        tb = _add_textbox(slide, 1.25, 1.65, 7.5, 1.0)
        _add_para(tb.text_frame, title, size=32, bold=True, color=_c("dk2"),
                  align=PP_ALIGN.CENTER, spacing=8)
        box = _add_filled_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                                1.25, 2.9, 7.5, 1.6, _c("lt1"))
        box.line.color.rgb = _c("accent1")
        tb = _add_textbox(slide, 1.45, 3.1, 7.1, 1.2)
        for line in info.split("\n"):
            _add_para(tb.text_frame, line, size=12, color=_c("dk2"),
                      align=PP_ALIGN.CENTER, spacing=6)
        _add_page_footer(slide, self._page, "")
        return slide

    def build(self, data):
        for s in data.get("slides", []):
            typ = s.get("type", "content")
            try:
                if typ == "cover":
                    self.build_cover(s)
                elif typ == "section":
                    self.build_section(s)
                elif typ == "agenda":
                    self.build_agenda(s)
                elif typ == "content":
                    self.build_content(s)
                elif typ == "two_column":
                    self.build_two_column(s)
                elif typ == "table":
                    self.build_table(s)
                elif typ == "image":
                    self.build_image_slide(s)
                elif typ == "summary":
                    self.build_summary(s)
                elif typ == "timeline":
                    self.build_timeline(s)
                elif typ == "chart":
                    self.build_chart(s)
                elif typ == "quote":
                    self.build_quote(s)
                elif typ == "contact":
                    self.build_contact(s)
                else:
                    print(f"unknown type {typ}", file=sys.stderr)
            except Exception as e:
                print(f"error building {typ} '{s.get('title','')}': {e}", file=sys.stderr)
                raise
        return self


def main():
    json_path = sys.argv[1] if len(sys.argv) > 1 else "docs/ppt/slides.json"
    out_path = sys.argv[2] if len(sys.argv) > 2 else "docs/智能门锁项目_立项与技术汇报_模板版.pptx"
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    builder = DeckBuilder("14 (1).pptx")
    builder.build(data)
    builder.prs.save(out_path)
    print(f"已生成 {out_path} ({builder._page} 页)")


if __name__ == "__main__":
    main()
