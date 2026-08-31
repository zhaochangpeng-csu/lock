#!/usr/bin/env python3
"""Beautify the 4 two-panel slides of 智能门锁.pptx with the 14(1).pptx theme.

Each slide is a two-column layout: left topic + right topic, rebuilt as
structured native cards/tables in the light-blue minimal style.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

SRC = "C:/Users/hoyo/OneDrive/文档/智能门锁.pptx"
OUT = "C:/Users/hoyo/OneDrive/文档/智能门锁_4页美化.pptx"
BG = "C:/Users/hoyo/Desktop/lock/docs/ppt/_tpl_bg/p01_0b91503e.jpg"

SW, SH = 13.333, 7.5

C_DK = RGBColor(0x44, 0x54, 0x6A)
C_A1 = RGBColor(0xAD, 0xC7, 0xDD)
C_A2 = RGBColor(0xE1, 0xF1, 0xFE)
C_A4 = RGBColor(0x6A, 0x9A, 0xC4)
C_A5 = RGBColor(0xA3, 0xC9, 0xE9)
C_RED = RGBColor(0xE0, 0x62, 0x5C)
C_GREEN = RGBColor(0x4F, 0xAE, 0x7A)
C_GRAY = RGBColor(0xA6, 0xA6, 0xA6)
C_AMBER = RGBColor(0xE5, 0xA6, 0x4F)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _no_shadow(sh):
    try:
        sh.shadow.inherit = False
    except Exception:
        pass


def _set(sh, lines, size, color, bold=False, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, font="DengXian", spacing=2):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = bold
            r.font.name = font; r.font.color.rgb = color


def add_bg(slide):
    slide.shapes.add_picture(BG, 0, 0, Inches(SW), Inches(SH))


def add_header(slide, chapter, title):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.6), Inches(0.42), Inches(3.0), Inches(0.42))
    pill.fill.solid(); pill.fill.fore_color.rgb = C_A4
    pill.line.fill.background(); _no_shadow(pill)
    _set(pill, chapter, 12, C_WHITE, bold=True)
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.94), Inches(12.1), Inches(0.5))
    _set(t, title, 24, C_DK, bold=True, align=PP_ALIGN.LEFT)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.5), Inches(1.0), Inches(0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb = C_A4
    ln.line.fill.background(); _no_shadow(ln)


def add_panel(slide, x, y, w, h, title, *, accent=C_A4):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = C_A1; card.line.width = Pt(1.5)
    _no_shadow(card)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.18), Inches(y + 0.18), Inches(0.07), Inches(0.42))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background(); _no_shadow(bar)
    tt = slide.shapes.add_textbox(Inches(x + 0.34), Inches(y + 0.14), Inches(w - 0.5), Inches(0.5))
    _set(tt, title, 17, C_DK, bold=True, align=PP_ALIGN.LEFT)
    return card


def add_table(slide, x, y, w, headers, rows, *, col_ratios=None, font=10.5, header_font=11):
    nrows = len(rows) + 1; ncols = len(headers)
    gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(0.34 * nrows))
    tbl = gf.table
    try:
        tbl.first_row = True; tbl.horz_banding = False
    except Exception:
        pass
    total = sum(col_ratios) if col_ratios else ncols
    for ci in range(ncols):
        r = col_ratios[ci] if col_ratios else 1
        tbl.columns[ci].width = Emu(int(Inches(w) * r / total))
    for ci, h in enumerate(headers):
        c = tbl.cell(0, ci)
        c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = C_A4
        c.margin_left = Inches(0.05); c.margin_right = Inches(0.05)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(header_font); r.font.bold = True
                r.font.name = "DengXian"; r.font.color.rgb = C_WHITE
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.text = val
            c.margin_left = Inches(0.05); c.margin_right = Inches(0.05)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = C_A2 if ri % 2 == 0 else C_WHITE
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(font); r.font.name = "DengXian"; r.font.color.rgb = C_DK
    return tbl


def add_row(slide, x, y, w, title, body, *, color=C_A4):
    """A compact horizontal entry: dot + title + body on one/two lines."""
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.08), Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    dot.line.fill.background(); _no_shadow(dot)
    if body:
        t = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y), Inches(w - 0.2), Inches(0.28))
        _set(t, [title + "　" + body], 12, C_DK, bold=False, align=PP_ALIGN.LEFT)
    else:
        t = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y), Inches(w - 0.2), Inches(0.28))
        _set(t, title, 12, C_DK, align=PP_ALIGN.LEFT)


def add_step(slide, num, x, y, w, title, body, *, color=C_A4):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.42), Inches(0.42))
    c.fill.solid(); c.fill.fore_color.rgb = color
    c.line.fill.background(); _no_shadow(c)
    _set(c, str(num), 14, C_WHITE, bold=True)
    t = slide.shapes.add_textbox(Inches(x + 0.55), Inches(y - 0.04), Inches(w - 0.55), Inches(0.5))
    _set(t, [title + "　" + body], 12.5, C_DK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


# ── slides ────────────────────────────────────────────────
def p1(slide):
    add_header(slide, "关键技术预研与选型", "技术选型 与 弃用路线")
    # left: tech selection table
    add_panel(slide, 0.6, 1.75, 6.05, 5.35, "技术选型：Demo 验证 vs 量产方案")
    add_table(slide, 0.85, 2.35, 5.55,
              ["能力域", "Demo 验证", "量产方案"],
              [["深度估计", "Depth Anything / MiDaS", "3D 结构光优先"],
               ["人脸识别", "InsightFace / ArcFace", "3D 结构光 + 端侧比对"],
               ["活体检测", "MediaPipe 动作活体", "3D 活体 + 分级唤醒"],
               ["声纹识别", "SpeechBrain ECAPA", "声纹作第二因子"],
               ["主认证兜底", "规则引擎 + 融合", "半导体指纹 + 密码"],
               ["多模态融合", "规则引擎加权", "规则化，策略从简"],
               ["执行推理", "ONNX Runtime / TensorRT", "全自动锁体 + NPU"]],
              col_ratios=[1.1, 1.8, 1.7], font=9.5, header_font=10.5)
    # right: rejected routes
    add_panel(slide, 6.8, 1.75, 5.93, 5.35, "弃用技术路线及理由", accent=C_RED)
    items = [("常电视觉", "功耗高、续航差"), ("云端比对", "隐私与断网风险"),
             ("光学指纹", "干湿手识别弱"), ("复杂 AI 融合", "可解释性差 · 留 V2.0"),
             ("UWB / 声纹开锁", "可行性低 / 仅第二因子"), ("指静脉", "识别慢、体积大")]
    y = 2.45
    for i, (t, b) in enumerate(items):
        add_row(slide, 7.1, y, 5.3, t, b, color=C_RED)
        y += 0.74


def p2(slide):
    add_header(slide, "系统架构与实现", "多模态认证：流程与融合")
    # left: 5-step flow
    add_panel(slide, 0.6, 1.75, 6.05, 5.35, "硬件多模态认证流程")
    steps = [("红外检测", "人体靠近"), ("人脸识别", "InsightFace"),
             ("活体检测", "眨眼 + 转头"), ("声纹识别", "文本无关"),
             ("融合判定", "加权求和")]
    y = 2.5
    for i, (t, b) in enumerate(steps, 1):
        add_step(slide, i, 0.95, y, 5.4, t, b)
        if i < 5:
            a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.15), Inches(y + 0.44), Inches(0.025), Inches(0.22))
            a.fill.solid(); a.fill.fore_color.rgb = C_A1
            a.line.fill.background(); _no_shadow(a)
        y += 0.66
    f = slide.shapes.add_textbox(Inches(0.95), Inches(5.9), Inches(5.4), Inches(0.4))
    _set(f, "权重 0.45/0.25/0.20/0.10 · 阈值 0.78（见右）", 10.5, C_GRAY, align=PP_ALIGN.LEFT)
    # right: fusion weights
    add_panel(slide, 6.8, 1.75, 5.93, 5.35, "多模态融合权重与阈值")
    add_table(slide, 7.05, 2.35, 5.4,
              ["模块", "权重", "阈值参数"],
              [["人脸", "0.45", "min_score 0.65"],
               ["活体", "0.25", "0.60 · 眨眼 + 转头"],
               ["声纹", "0.20", "0.70 · 最短 2s"],
               ["红外", "0.10", "上升沿 · 保持 30s"],
               ["融合阈值", "0.78", "加权总分判定"]],
              col_ratios=[1.2, 0.9, 2.2], font=11, header_font=11)
    f2 = slide.shapes.add_textbox(Inches(7.05), Inches(5.6), Inches(5.4), Inches(0.4))
    _set(f2, "参数在 config.yaml 可调，Agent 不参与评分", 10.5, C_GRAY, align=PP_ALIGN.LEFT)


def p3(slide):
    add_header(slide, "系统架构与实现", "语音 Agent：流水线与受控工具")
    # left: pipeline 6 steps
    add_panel(slide, 0.6, 1.75, 6.05, 5.35, "低延迟语音 Agent 流水线")
    steps = [("麦克风", "16 kHz 采集"), ("VAD", "起止 / 打断"),
             ("ASR", "本地识别"), ("LLM", "流式回复"),
             ("TTS", "本地合成"), ("播放", "aplay 输出")]
    y = 2.5
    for i, (t, b) in enumerate(steps, 1):
        add_step(slide, i, 0.95, y, 5.4, t, b)
        if i < 6:
            a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.15), Inches(y + 0.44), Inches(0.025), Inches(0.18))
            a.fill.solid(); a.fill.fore_color.rgb = C_A1
            a.line.fill.background(); _no_shadow(a)
        y += 0.60
    f = slide.shapes.add_textbox(Inches(0.95), Inches(6.25), Inches(5.4), Inches(0.4))
    _set(f, "ASR / TTS 本地执行 · DeepSeek 需网络", 10.5, C_GRAY, align=PP_ALIGN.LEFT)
    # right: tools
    add_panel(slide, 6.8, 1.75, 5.93, 5.35, "Agent 受控工具（仅两个）")
    add_table(slide, 7.05, 2.35, 5.4,
              ["工具", "方法", "用途"],
              [["current_auth_context", "GET /tools/…", "读取最新凭证"],
               ["request_unlock", "POST /tools/…", "请求开锁"]],
              col_ratios=[1.7, 1.5, 1.3], font=11, header_font=11)
    for i, (t, b) in enumerate([
            ("开锁前必须先调用", "current_auth_context 四条件校验"),
            ("不暴露裸 unlock", "Agent 不能评分 / 否决 / 伪造")]):
        add_row(slide, 7.1, 4.1 + i * 0.75, 5.3, t, b, color=C_AMBER)
    f2 = slide.shapes.add_textbox(Inches(7.05), Inches(6.3), Inches(5.4), Inches(0.4))
    _set(f2, "安全闸门拥有最终裁决权", 10.5, C_RED, align=PP_ALIGN.LEFT)


def p4(slide):
    add_header(slide, "安全设计与异常闭环", "四层安全约束 与 场景验证")
    # left: 4 layers
    add_panel(slide, 0.6, 1.75, 6.05, 5.35, "四层安全约束")
    layers = [("第一层 · 融合层", "只有本地 FusionEngine 产生认证结果"),
              ("第二层 · 凭证层", "短时一次性，过期重复消费一律拒绝"),
              ("第三层 · Agent 层", "仅两个受控工具，不暴露裸 unlock"),
              ("第四层 · 执行层", "Bearer Token + 干运行 + 本地执行器")]
    y = 2.45
    for t, b in layers:
        add_row(slide, 1.0, y, 5.4, t, b, color=C_A4)
        y += 0.74
    # conclusion band
    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.5), Inches(5.4), Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = C_A2
    band.line.color.rgb = C_A4; band.line.width = Pt(1.5); _no_shadow(band)
    _set(band, "即使模型误解，本地闸门仍拒绝无凭证请求", 13, C_DK, bold=True)
    # right: 2 scenarios
    add_panel(slide, 6.8, 1.75, 5.93, 5.35, "安全验证：两个真实场景")
    sc = [("有效凭证 + 「请开门」", "current_auth_context → request_unlock", "干运行阻止真实动作", C_GREEN),
          ("无效凭证 + 「请开门」", "仅 current_auth_context", "不调用开锁工具", C_RED)]
    for i, (t, path, res, col) in enumerate(sc):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(2.45 + i * 1.85), Inches(5.3), Inches(1.6))
        card.fill.solid(); card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = col; card.line.width = Pt(1.5); _no_shadow(card)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.1), Inches(2.55 + i * 1.85), Inches(0.07), Inches(1.4))
        bar.fill.solid(); bar.fill.fore_color.rgb = col
        bar.line.fill.background(); _no_shadow(bar)
        tt = slide.shapes.add_textbox(Inches(7.35), Inches(2.55 + i * 1.85), Inches(4.9), Inches(0.4))
        _set(tt, t, 14, C_DK, bold=True, align=PP_ALIGN.LEFT)
        bb = slide.shapes.add_textbox(Inches(7.35), Inches(2.95 + i * 1.85), Inches(4.9), Inches(0.9))
        _set(bb, [path, "结果：" + res], 11.5, C_DK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    f2 = slide.shapes.add_textbox(Inches(7.1), Inches(6.3), Inches(5.3), Inches(0.4))
    _set(f2, "结论：模型无法绕过本地安全闸门", 12, C_DK, bold=True, align=PP_ALIGN.LEFT)


def main():
    prs = Presentation(SRC)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[1]
    for fn in (p1, p2, p3, p4):
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        fn(slide)
    prs.save(OUT)
    print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
