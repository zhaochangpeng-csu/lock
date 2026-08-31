# -*- coding: utf-8 -*-
"""体检候选 PPT 模板：尺寸、页数、版式、占位符、主题色、动画。"""
import sys
from pptx import Presentation
from pptx.util import Emu

FILES = [
    r"C:\Users\hoyo\Desktop\lock\PPT模板.pptx",
    r"C:\Users\hoyo\Desktop\lock\14 (1).pptx",
    r"C:\Users\hoyo\Desktop\lock\16 (6).pptx",
]

PH_TYPE = {
    0: "TITLE", 1: "BODY", 2: "CENTER_TITLE", 3: "SUBTITLE", 4: "OTHER",
    5: "VERT_BODY", 6: "OBJECT", 7: "VERT_TITLE", 8: "BITMAP", 13: "PICTURE",
    14: "TABLE", 15: "CLIP_ART", 16: "SLIDE_NUM", 17: "DATE", 18: "FOOTER",
}


def emu_in(v):
    return round(Emu(v).inches, 2) if v is not None else None


for path in FILES:
    print("=" * 78)
    print("FILE:", path)
    try:
        prs = Presentation(path)
    except Exception as e:
        print("  OPEN FAILED:", e)
        continue

    w, h = prs.slide_width, prs.slide_height
    ratio = round(Emu(w).inches / Emu(h).inches, 3)
    print(f"  size: {emu_in(w)} x {emu_in(h)} in   ratio={ratio}")
    print(f"  slides: {len(prs.slides)}   layouts: {len(prs.slide_layouts)}")

    # 主题色
    try:
        scheme = prs.slide_masters[0].slide_layouts[0]
        theme = prs.slide_masters[0].element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme"
        )
        if theme is not None:
            names = []
            for child in theme:
                tag = child.tag.split("}")[-1]
                if tag in ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2",
                           "accent3", "accent4", "accent5", "accent6",
                           "hlink", "folHlink"):
                    val = child.get("val") or (
                        child[0].get("val") if len(child) else "?")
                    names.append(f"{tag}=#{val}")
            print("  theme:", ", ".join(names[:8]))
    except Exception as e:
        print("  theme: (read fail)", e)

    # 各版式
    print("  -- layouts --")
    for i, lay in enumerate(prs.slide_layouts):
        phs = []
        for p in lay.placeholders:
            t = PH_TYPE.get(p.placeholder_format.type, str(p.placeholder_format.type))
            phs.append(f"{p.placeholder_format.idx}:{t}")
        shape_n = len(lay.shapes)
        print(f"   [{i:>2}] {lay.name:<28} shapes={shape_n:<3} ph={','.join(phs)}")

    # 动画统计
    print("  -- animations --")
    total_anim = 0
    total_morph = 0
    for i, s in enumerate(prs.slides):
        timing = s.element.find(
            ".//{http://schemas.openxmlformats.org/presentationml/2006/main}timing")
        n_anim = 0
        if timing is not None:
            n_anim = timing.findall(
                ".//{http://schemas.openxmlformats.org/presentationml/2006/main}animEffect") \
                .__len__()
        trans = s.element.find(
            ".//{http://schemas.openxmlformats.org/presentationml/2006/main}transition")
        is_morph = False
        if trans is not None:
            for ch in trans:
                if "morph" in ch.tag.lower():
                    is_morph = True
        total_anim += n_anim
        total_morph += 1 if is_morph else 0
    print(f"  animEffect total={total_anim}  morph-transition slides={total_morph}")

    # 前几页内容预览
    print("  -- first slides text --")
    for i, s in enumerate(prs.slides):
        if i >= 4:
            break
        texts = []
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                texts.append(sh.text_frame.text.strip().replace("\n", " / ")[:60])
        print(f"   slide{i}: {texts[:4]}")
    print()
