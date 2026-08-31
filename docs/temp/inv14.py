# -*- coding: utf-8 -*-
"""对 14 (1).pptx 做逐页逐形状清点，为内容对位做准备。"""
from pptx import Presentation
from pptx.util import Emu

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

PATH = r"C:\Users\hoyo\Desktop\lock\14 (1).pptx"
prs = Presentation(PATH)
SW = Emu(prs.slide_width).inches
SH = Emu(prs.slide_height).inches
print(f"canvas: {SW} x {SH} in, slides={len(prs.slides)}\n")


def pos(sh):
    try:
        return (round(Emu(sh.left).inches, 2), round(Emu(sh.top).inches, 2),
                round(Emu(sh.width).inches, 2), round(Emu(sh.height).inches, 2))
    except Exception:
        return None


for i, s in enumerate(prs.slides):
    lay = s.slide_layout
    print(f"──── slide {i}  layout=[{lay.name}]")
    for sh in s.shapes:
        kind = sh.shape_type
        p = pos(sh)
        txt = ""
        if sh.has_text_frame:
            t = sh.text_frame.text.replace("\n", "⏎").strip()
            txt = t[:70]
        npar = len(sh.text_frame.paragraphs) if sh.has_text_frame else 0
        tag = "PIC" if sh.shape_type == 13 else (
            "GRP" if sh.shape_type == 6 else (
                "TBL" if sh.shape_type == 19 else "SP"))
        extra = ""
        if tag == "PIC":
            extra = f" img={getattr(sh, 'image', None) and sh.image.filename[:22]}"
        print(f"   {tag:<4} id={sh.shape_id:<4} name='{sh.name[:26]}' "
              f"pos={p} par={npar}{extra}")
        if txt:
            print(f"        text: {txt}")
    print()
