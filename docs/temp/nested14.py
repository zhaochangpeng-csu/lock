# -*- coding: utf-8 -*-
"""递归列出每页所有文本框（含组合内部），按文档顺序，供内容对位。"""
import sys
from pptx import Presentation
from pptx.util import Emu

PATH = r"C:\Users\hoyo\Desktop\lock\14 (1).pptx"
prs = Presentation(PATH)


def walk(shapes, depth, out, prefix=""):
    for sh in shapes:
        p = None
        try:
            p = (round(Emu(sh.left).inches, 2), round(Emu(sh.top).inches, 2),
                 round(Emu(sh.width).inches, 2), round(Emu(sh.height).inches, 2))
        except Exception:
            pass
        is_grp = sh.shape_type == 6
        if sh.has_text_frame:
            t = sh.text_frame.text.replace("\n", " ⏎ ").strip()
            out.append((depth, f"{prefix}{sh.name}", p, t[:78], len(sh.text_frame.paragraphs)))
        if is_grp:
            try:
                walk(sh.shapes, depth + 1, out, prefix + "  ")
            except Exception:
                pass


only = [int(x) for x in sys.argv[1:]] if len(sys.argv) > 1 else None

for i, s in enumerate(prs.slides):
    if only is not None and i not in only:
        continue
    out = []
    walk(s.shapes, 0, out)
    print(f"═══ slide {i}  layout=[{s.slide_layout.name}]  textframes={len(out)}")
    for d, nm, p, t, npar in out:
        print(f"  {'·'*(d+1)} [{npar}p] {nm[:34]:<34} {str(p):<28} {t}")
    print()
