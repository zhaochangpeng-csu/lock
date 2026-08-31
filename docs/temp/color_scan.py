# -*- coding: utf-8 -*-
"""扫描每个模板的主色（按面积加权近似：统计 solidFill 出现频次与图形类型）。"""
from collections import Counter
from pptx import Presentation

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

FILES = [
    ("PPT模板.pptx", r"C:\Users\hoyo\Desktop\lock\PPT模板.pptx"),
    ("14 (1).pptx", r"C:\Users\hoyo\Desktop\lock\14 (1).pptx"),
    ("16 (6).pptx", r"C:\Users\hoyo\Desktop\lock\16 (6).pptx"),
]


def lum(h):
    if not h or len(h) != 6:
        return None
    try:
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    except ValueError:
        return None
    return (r * 299 + g * 587 + b * 114) / 1000


for name, path in FILES:
    prs = Presentation(path)
    fills = Counter()
    for s in prs.slides:
        for el in s.element.iter():
            if el.tag == f"{{{NS_A}}}srgbClr":
                fills[el.get("val")] += 1
    # 主题色映射
    master = prs.slide_masters[0]
    theme_el = master.element.find(f".//{{{NS_A}}}clrScheme")
    pal = {}
    if theme_el is not None:
        for ch in theme_el:
            tag = ch.tag.split("}")[-1]
            c = ch.find(f"{{{NS_A}}}srgbClr")
            if c is not None:
                pal[tag] = c.get("val")

    top = fills.most_common(10)
    print(f"### {name}  (slides={len(prs.slides)})")
    print("  theme:", {k: "#" + v for k, v in pal.items()
                       if k in ("dk1", "lt1", "dk2", "lt2", "accent1",
                                "accent2", "accent3", "accent4", "accent5", "accent6")})
    print("  top colors used:")
    for h, c in top:
        L = lum(h)
        tone = "dark" if (L is not None and L < 80) else (
            "light" if (L is not None and L > 190) else "mid")
        print(f"    #{h}  x{c:<4} lum={round(L,1) if L else '?':>6}  {tone}")
    dark_share = sum(c for h, c in fills.items()
                     if (lum(h) or 255) < 80)
    tot = sum(fills.values())
    print(f"  => dark-color share: {round(dark_share / tot * 100, 1)}% "
          f"(of {tot} solid fills)")
    print()
