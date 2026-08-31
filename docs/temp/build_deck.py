# -*- coding: utf-8 -*-
"""
基于 14 (1).pptx 模板生成 17 页初稿。
核心原则：复制模板原页 -> 保留全部形状与动画(含 p:timing) -> 按原文定位替换文案。
绝不删除动画宿主形状。
"""
import copy
import os
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

import os as _os
TPL = r"C:\Users\hoyo\Desktop\lock\14 (1).pptx"
ARCH = r"C:\Users\hoyo\Desktop\lock\docs\architecture.png"
UNLOCK = r"C:\Users\hoyo\Desktop\lock\docs\unlock_sequence.png"
OUT = _os.environ.get(
    "PPT_OUT", r"C:\Users\hoyo\Desktop\lock\docs\智能门锁_语音Agent_初稿.pptx")
COPY_TIMING = _os.environ.get("PPT_NO_TIMING") != "1"
COPY_TRANS = _os.environ.get("PPT_NO_TRANS") != "1"
COPY_SHAPES = _os.environ.get("PPT_NO_SHAPES") != "1"
COPY_RELS = _os.environ.get("PPT_NO_RELS") != "1"
SKIP_RELS = [k for k in _os.environ.get("PPT_SKIP_RELS", "").split(",") if k]

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

prs = Presentation(TPL)
SW = Emu(prs.slide_width).inches
SH = Emu(prs.slide_height).inches
print(f"template canvas {SW} x {SH}  slides={len(prs.slides)}")

FONT = "微软雅黑"
INK = RGBColor(0x33, 0x33, 0x33)


# ─────────────────────────── 底层工具 ───────────────────────────

def collect_tfs(slide):
    """按文档顺序收集所有文本框（含组合内部）。"""
    out = []

    def walk(shapes):
        for sh in shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame)
            if sh.shape_type == 6:  # GROUP
                try:
                    walk(sh.shapes)
                except Exception:
                    pass
    walk(slide.shapes)
    return out


def collect_text_shapes(slide):
    """按文档顺序收集所有带文本框的形状（含组合内部）。"""
    out = []

    def walk(shapes):
        for sh in shapes:
            if sh.has_text_frame:
                out.append(sh)
            if sh.shape_type == 6:
                try:
                    walk(sh.shapes)
                except Exception:
                    pass
    walk(slide.shapes)
    return out


def widen(slide, orig, w=None, h=None, nth=0):
    """按原文定位形状并放大（组合内使用内部坐标，宽高单位为英寸）。"""
    cands = [sh for sh in collect_text_shapes(slide)
             if sh.text_frame.text.strip() == orig]
    if not cands:
        cands = [sh for sh in collect_text_shapes(slide)
                 if orig in sh.text_frame.text]
    if len(cands) <= nth:
        print(f"    [miss] widen: {orig[:24]!r} nth={nth}")
        return False
    sh = cands[nth]
    if w is not None:
        sh.width = Inches(w)
    if h is not None:
        sh.height = Inches(h)
    return True


def _ensure_run(p_el):
    """确保段落有一个 run，返回该 run 的 a:t 元素。"""
    runs = p_el.findall(qn("a:r"))
    if runs:
        for r in runs[1:]:
            p_el.remove(r)
        r = runs[0]
    else:
        r = p_el.makeelement(qn("a:r"), {})
        p_el.append(r)
    t = r.find(qn("a:t"))
    if t is None:
        t = r.makeelement(qn("a:t"), {})
        r.append(t)
    return t


def set_tf(tf, lines, size=None, bold=None, color=None, align=None):
    """写入文本，保留模板原有字体格式（复用首个 run 的 rPr）。"""
    if isinstance(lines, str):
        lines = [lines]
    lines = list(lines)
    body = tf._txBody
    paras = body.findall(qn("a:p"))
    if not paras:
        return False
    base = paras[0]
    for p in paras[1:]:
        body.remove(p)
    for i, ln in enumerate(lines):
        el = base if i == 0 else copy.deepcopy(base)
        for br in el.findall(qn("a:br")):
            el.remove(br)
        if i > 0:
            body.append(el)
        t = _ensure_run(el)
        t.text = ln
        if size is not None or bold is not None or color is not None:
            runs = el.findall(qn("a:r"))
            if runs:
                rpr = runs[0].find(qn("a:rPr"))
                if rpr is None:
                    rpr = runs[0].makeelement(qn("a:rPr"), {})
                    runs[0].insert(0, rpr)
                if size is not None:
                    rpr.set("sz", str(int(size * 100)))
                if bold is not None:
                    rpr.set("b", "1" if bold else "0")
                if color is not None:
                    for c in rpr.findall(qn("a:solidFill")):
                        rpr.remove(c)
                    sf = rpr.makeelement(qn("a:solidFill"), {})
                    clr = sf.makeelement(qn("a:srgbClr"), {})
                    clr.set("val", str(color))
                    sf.append(clr)
                    rpr.append(sf)
                # 中文字体
                latin = rpr.find(qn("a:latin"))
                if latin is None:
                    latin = rpr.makeelement(qn("a:latin"), {})
                    rpr.append(latin)
                latin.set("typeface", FONT)
    return True


def set_by_orig(slide, orig, new, nth=0, size=None, bold=None, color=None):
    """按模板原文定位文本框并替换。返回是否命中。"""
    tfs = collect_tfs(slide)
    hits = [tf for tf in tfs if tf.text.strip() == orig]
    if not hits:
        hits = [tf for tf in tfs if orig in tf.text]
    if len(hits) <= nth:
        print(f"    [miss] orig={orig[:26]!r} nth={nth} (hits={len(hits)})")
        return False
    return set_tf(hits[nth], new, size=size, bold=bold, color=color)


def chip_group(slide):
    """模板各内容页共用的页眉标题条，固定位于 (4.44, 0.2)。"""
    for sh in slide.shapes:
        try:
            if (abs(Emu(sh.left).inches - 4.44) < 0.25 and
                    abs(Emu(sh.top).inches - 0.2) < 0.25):
                return sh
        except Exception:
            pass
    return None


def set_chip(slide, zh, en=None, size=None):
    """按位置定位标题条（不能用原文定位：它在文档顺序里排在卡片之后）。"""
    g = chip_group(slide)
    if g is None:
        print("    [miss] no title chip on this slide")
        return False
    tfs = []

    def walk(shapes):
        for sh in shapes:
            if sh.has_text_frame:
                tfs.append(sh.text_frame)
            if sh.shape_type == 6:
                try:
                    walk(sh.shapes)
                except Exception:
                    pass
    walk(g.shapes)
    if not tfs:
        return False
    set_tf(tfs[0], zh, size=size or 14, bold=True)
    if en and len(tfs) > 1:
        set_tf(tfs[1], en, size=(size or 14) - 4)
    return True


def duplicate(src_idx):
    """复制模板第 src_idx 页，保留形状、图片关系与 p:timing 动画。"""
    src = prs.slides[src_idx]
    layout = src.slide_layout
    dest = prs.slides.add_slide(layout)
    # 清掉 add_slide 自动带的占位符
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)
    # 复制形状
    if COPY_SHAPES:
        for shp in src.shapes:
            dest.shapes._spTree.append(copy.deepcopy(shp._element))
    # 复制关系（图片/图表/媒体）；rId 冲突时分配新 id 并重映射引用
    remap = {}
    for rId, rel in (src.part.rels.items() if COPY_RELS else []):
        if "notesSlide" in rel.reltype:
            continue
        if rel.reltype.endswith("/slideLayout"):
            continue          # add_slide 已经建好，重复会产生双重版式关系
        if rel.is_external:
            continue
        # themeOverride 必须跳过：add_slide 已建立版式/主题关系，
        # 再复制来源页的 themeOverride 会让 PowerPoint 判定文件损坏。
        if rel.reltype.endswith("/themeOverride"):
            continue
        if any(k in rel.reltype for k in SKIP_RELS):
            continue
        try:
            new = dest.part.rels.get_or_add(rel.reltype, rel.target_part)
        except Exception as e:
            print(f"    [rel warn] {rId}: {e}")
            continue
        if new != rId:
            remap[rId] = new
    if remap:
        for el in dest._element.iter():
            for attr in (qn("r:id"), qn("r:embed"), qn("r:link")):
                v = el.get(attr)
                if v in remap:
                    el.set(attr, remap[v])
    # 复制动画：p:timing 与 p:transition
    dest_el = dest._element
    for tag in ("transition", "timing"):
        if tag == "transition" and not COPY_TRANS:
            continue
        if tag == "timing" and not COPY_TIMING:
            continue
        node = src._element.find(qn(f"p:{tag}"))
        if node is None:
            continue
        old = dest_el.find(qn(f"p:{tag}"))
        if old is not None:
            dest_el.remove(old)
        ext = dest_el.find(qn("p:extLst"))
        if ext is not None:
            ext.addprevious(copy.deepcopy(node))
        else:
            dest_el.append(copy.deepcopy(node))
    return dest


def delete_slides(indices):
    """按索引删除幻灯片（含关系清理）。"""
    lst = prs.slides._sldIdLst
    for idx in sorted(indices, reverse=True):
        sldId = lst[idx]
        rId = sldId.get(qn("r:id"))
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        lst.remove(sldId)


def find_shape(slide, name):
    out = []

    def walk(shapes):
        for sh in shapes:
            if sh.name == name:
                out.append(sh)
            if sh.shape_type == 6:
                try:
                    walk(sh.shapes)
                except Exception:
                    pass
    walk(slide.shapes)
    return out[0] if out else None


def add_text(slide, l, t, w, h, lines, size=12, bold=False,
             color=INK, align=None, spacing=1.25, after=3):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(ln, tuple):
            txt, opts = ln
        else:
            txt, opts = ln, {}
        p.text = txt
        p.line_spacing = opts.get("spacing", spacing)
        p.space_after = Pt(opts.get("after", after))
        for r in p.runs:
            r.font.size = Pt(opts.get("size", size))
            r.font.bold = opts.get("bold", bold)
            r.font.color.rgb = opts.get("color", color)
            r.font.name = FONT
        if opts.get("align"):
            p.alignment = opts["align"]
    return tb


def replace_pic(slide, shape, img_path, keep_aspect=True):
    """替换图片内容（保留形状本身，动画不丢）。"""
    if not os.path.exists(img_path):
        print(f"    [miss] image {img_path}")
        return False
    el = shape._element
    blips = el.findall(".//" + qn("a:blip"))
    if not blips:
        return False
    # 借 add_picture 建好图片关系，再把 rId 挪到目标形状上，最后删掉临时图
    try:
        tmp = slide.shapes.add_picture(img_path, Inches(0), Inches(0))
        iw, ih = tmp.image.size
        rId = tmp._element.find(".//" + qn("a:blip")).get(qn("r:embed"))
        blips[0].set(qn("r:embed"), rId)
        tmp._element.getparent().remove(tmp._element)
    except Exception as e:
        print(f"    [err] add image: {e}")
        return False
    if keep_aspect:
        box_w = Emu(shape.width).inches
        box_h = Emu(shape.height).inches
        ar_img = iw / ih
        ar_box = box_w / box_h
        if ar_img > ar_box:      # 图更宽 -> 以宽为准
            nh = box_w / ar_img
            shape.width = Inches(box_w)
            shape.height = Inches(nh)
            shape.top = Inches(Emu(shape.top).inches + (box_h - nh) / 2)
        else:                    # 图更高 -> 以高为准
            nw = box_h * ar_img
            shape.height = Inches(box_h)
            shape.width = Inches(nw)
            shape.left = Inches(Emu(shape.left).inches + (box_w - nw) / 2)
    return True


def add_table(slide, l, t, w, h, rows, col_w, size=10.5, head_size=11):
    nr, nc = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nr, nc, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    tbl = gf.table
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        tbl.rows[i].height = Inches(h / nr)
        for j, cell in enumerate(row):
            tf = tbl.cell(i, j).text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.06)
            tf.margin_top = tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]
            p.text = cell
            for r in p.runs:
                r.font.size = Pt(head_size if i == 0 else size)
                r.font.bold = (i == 0)
                r.font.name = FONT
                r.font.color.rgb = (RGBColor(0xFF, 0xFF, 0xFF) if i == 0
                                    else INK)
    return gf


# ─────────────────────────── 17 页内容 ───────────────────────────

built = []


def page(tag):
    print(f"\n── {tag}")
    return None


# 01 封面 ────────────────────────────────────────────
page("01 封面  <- s0")
s = duplicate(0)
set_by_orig(s, "极简淡蓝", "基于 Jetson 的多模态智能门锁", size=30, bold=True)
set_by_orig(s, "商务通用工作总结ppt模板",
            "硬件认证 · 一次性凭证 · 语音交互 · 事件闭环", size=13)
set_by_orig(s, "Business Work Summary Report Template",
            "SMART LOCK × VOICE AGENT  ·  MVP DEMO 2026", size=10)
set_by_orig(s, "汇报人", "", size=9)          # 模板隐藏署名位，清空
set_by_orig(s, "XXXXL", "", size=9)
built.append(s)

# 02 目录 ────────────────────────────────────────────
page("02 目录  <- s1")
s = duplicate(1)
TOC = [
    ("01  项目背景", "重新审视门锁的三个缺口，明确四个目标与边界"),
    ("02  技术方案", "边缘融合认证、两阶段开门与低延迟语音流水线"),
    ("03  系统实现", "四层安全约束、异常事件闭环与部署进程治理"),
    ("04  验证与规划", "测试结果、已知边界与下一步计划"),
]
for i, (t_, d_) in enumerate(TOC):
    set_by_orig(s, "请输入你的题目", t_, nth=0, size=13, bold=True)
    set_by_orig(s, "请在此处添加详细描述文本，尽量与标题文本语言风格相符合，"
                   "语言描述尽量简洁生动。", d_, nth=0, size=9)
built.append(s)

# 03 为什么重新设计门锁 ────────────────────────────────
page("03 项目背景：为什么重新设计门锁  <- s8")
s = duplicate(8)
set_chip(s, "项目背景", "BACKGROUND")
t = find_shape(s, "矩形 106")          # 居中主标题，原框过窄需加宽
if t is not None:
    t.width = Inches(4.2)
    t.left = Inches(5.0 - 2.1)
set_by_orig(s, "输入标题", "为什么重新设计门锁", nth=0, size=15, bold=True)
set_by_orig(s, "We are more powerful when we are clear",
            "单一模态撑不起门锁安全，自然交互与异常记录同样缺位。", size=9.5)
set_by_orig(s, "Clarity Leads to Power.", "", size=9)
PAIN = [
    ("单一识别易被冒用", "照片、替身与环境噪声都能绕过单一模态，人脸或声纹单独使用都存在明确攻击面。"),
    ("缺少自然交互", "传统门锁只有按键与卡片，没有语音、没有对话，用户意图无法被理解。"),
    ("异常缺少统一记录", "陌生人徘徊、认证失败不落盘，也没有后续通知与追溯入口。"),
    ("大模型不宜直接控锁", "语言模型输出存在不确定性，不能承担开锁这类安全判定职责。"),
]
for i, (t_, d_) in enumerate(PAIN):
    set_by_orig(s, "Header here please", t_, nth=0, size=12, bold=True)
    set_by_orig(s, "We are placing great text in this text frames.", d_,
                nth=0, size=9)
built.append(s)

# 04 四个目标，四条边界 ────────────────────────────────
page("04 四个目标，四条边界  <- s15")
s = duplicate(15)
set_chip(s, "四个目标，四条边界", "OBJECTIVES")
set_by_orig(s, "Ipsum dolor sit amet",
            "对话能力增强体验，硬件认证保持安全权威。", size=12)
GOALS = [
    ("边缘端融合认证", "红外、人脸、活体、声纹四路在 Jetson 本地完成"),
    ("自然语音开门", "认证通过后自然语言表达意图，无需固定口令"),
    ("受控工具调用", "白名单工具与令牌鉴权，动作由本地闸门决定"),
    ("异常自动落盘", "认证失败写入事件文件，为通知与追溯提供入口"),
]
for i in range(4):      # 卡片文本框原尺寸过小，先放大再填字
    widen(s, "输入标题", w=2.5, h=0.5, nth=i)
    widen(s, "Add your title here pleases", w=2.5, h=1.15, nth=i)
for i, (t_, d_) in enumerate(GOALS):
    set_by_orig(s, "输入标题", t_, nth=0, size=11, bold=True)
    set_by_orig(s, "Add your title here pleases", d_, nth=0, size=8.5)
for i, num in enumerate(["01", "02", "03", "04"]):
    set_by_orig(s, ["1994", "07", "04", "356"][i], num, size=20, bold=True)
built.append(s)

# 05 总体部署架构 ─────────────────────────────────────
page("05 总体部署架构  <- s11")
s = duplicate(11)
set_chip(s, "总体部署架构", "ARCHITECTURE")
sq = find_shape(s, "矩形 2")          # 左侧方形图区 (0.81,0.88,3.86,3.86)
if sq is not None:
    pic = s.shapes.add_picture(ARCH, Inches(0.81), Inches(0.88),
                               Inches(3.86))
    ar = pic.image.size[0] / pic.image.size[1]
    nh = 3.86 / ar
    pic.height = Inches(nh)
    pic.top = Inches(0.88 + (3.86 - nh) / 2)
set_by_orig(s, "单击此处输入标题文本", "Jetson 边缘端", nth=0, size=13, bold=True)
set_by_orig(s, "I love you more than", "硬件采集、多模态识别与融合判定；本地 ASR/TTS 与 Pipecat 流水线；一次性凭证、工具网关、异常服务与最终执行器。",
            nth=0, size=9)
set_by_orig(s, "单击此处输入标题文本", "Windows / WSL", nth=0, size=13, bold=True)
set_by_orig(s, "I love you more than", "FastGPT 工作流、DeepSeek 调用、SSH 隧道与部署管理。DeepSeek 不接触摄像头特征、声纹 embedding 或真实执行器。",
            nth=0, size=9)
built.append(s)

# 06 硬件多模态认证 ────────────────────────────────────
page("06 硬件多模态认证  <- s18")
s = duplicate(18)
set_chip(s, "硬件多模态认证", "MULTIMODAL AUTH")
chart = find_shape(s, "Chart 4")
if chart is not None:
    chart._element.getparent().remove(chart._element)   # 未被动画引用
add_text(s, 0.82, 1.55, 4.26, 0.4, "认证链路", size=13, bold=True)
add_text(s, 0.82, 2.0, 4.26, 2.3, [
    "红外触发 → InsightFace 人脸识别",
    "→ MediaPipe 眨眼 / 转头活体",
    "→ SpeechBrain ECAPA 声纹（文本无关）",
    "→ 人脸 / 声纹身份一致性校验",
    "→ 本地加权融合判定",
    "融合阈值 0.78，默认 require_all=false",
    "声纹阈值 0.70，最短有效语音 2 秒",
], size=10.5, spacing=1.5, after=5)
STEPS = [
    ("红外触发", "人体靠近，唤醒摄像头"),
    ("人脸识别", "InsightFace / ArcFace 身份判定"),
    ("活体检测", "眨眼一次 + 左右转头"),
    ("声纹识别", "自然说话，无需固定口令"),
    ("融合判定", "加权总分 ≥ 0.78 视为通过"),
]
for i, (t_, d_) in enumerate(STEPS):
    set_by_orig(s, ["One Category", "Two Category", "Three Category",
                    "Four Category", "Five Category"][i], t_, size=11,
                bold=True)
    set_by_orig(s, "When an unknown printer took a galley of type", d_,
                nth=0, size=8.5)
built.append(s)

# 07 两阶段受控开门 ────────────────────────────────────
page("07 两阶段受控开门  <- s3")
s = duplicate(3)
pic = find_shape(s, "图片 1")
if pic is not None:
    replace_pic(s, pic, UNLOCK)          # 保留形状 = 保留动画
add_text(s, 0.7, 2.72, 6.0, 0.42, "两阶段受控开门", size=15, bold=True)
set_by_orig(s, "关于我们", ["阶段一 · 硬件认证", "PHASE 1 · HARDWARE"],
            size=11, bold=True)
set_by_orig(s, "标题数字等都可以通过点击和重新输入进行更改",
            "融合通过后不立即开锁，Jetson 写入一次性凭证：含签发时间、fusion_passed 与 consumed，有效期 300 秒，只能消费一次。"
            "用户说「开门」后，Agent 必须先调用 current_auth_context，仅当凭证可用、新鲜、已授权且未消费时才调用 request_unlock；"
            "本地网关再次校验后执行。lock.flow=immediate 可恢复原直开流程。", size=9.5)
built.append(s)

# 08 低延迟语音 Agent 流水线 ────────────────────────────
page("08 语音 Agent 流水线  <- s18")
s = duplicate(18)
set_chip(s, "低延迟语音 Agent 流水线", "VOICE PIPELINE")
chart = find_shape(s, "Chart 4")
if chart is not None:
    chart._element.getparent().remove(chart._element)
add_text(s, 0.82, 1.55, 4.26, 0.4, "音频链路", size=13, bold=True)
add_text(s, 0.82, 2.0, 4.26, 2.3, [
    "麦克风：16 kHz 单声道，PortAudio",
    "Pipecat + Silero VAD：轮流说话与打断",
    "本地 ASR：SenseVoice / FunASR",
    "Agent：FastGPT 工作流 + DeepSeek",
    "本地 TTS：sherpa-onnx 优先",
    "输出：ALSA aplay，16 kHz S16_LE",
    "",
], size=10.5, spacing=1.5, after=5)
PIPE = [("采集", "sounddevice / PortAudio 16kHz"),
        ("分段", "Silero VAD 与打断处理"),
        ("识别", "本地 SenseVoice / FunASR"),
        ("决策", "FastGPT + DeepSeek"),
        ("合成", "sherpa-onnx → aplay")]
for i, (t_, d_) in enumerate(PIPE):
    set_by_orig(s, ["One Category", "Two Category", "Three Category",
                    "Four Category", "Five Category"][i], t_, size=11,
                bold=True)
    set_by_orig(s, "When an unknown printer took a galley of type", d_,
                nth=0, size=8.5)
built.append(s)

# 09 四层安全约束 ──────────────────────────────────────
page("09 四层安全约束  <- s15")
s = duplicate(15)
set_chip(s, "四层安全约束", "SECURITY LAYERS")
set_by_orig(s, "Ipsum dolor sit amet",
            "即使模型误解或提示词失效，本地安全闸门仍能拒绝无凭证请求。", size=12)
SEC = [("融合层", "只有 Jetson 本地 FusionEngine 能产出认证结果"),
       ("凭证层", "短时、一次性，过期与重复消费均拒绝"),
       ("Agent 层", "只注册 current_auth_context 与 request_unlock"),
       ("执行层", "Bearer Token + NO_UNLOCK 双重闸门")]
for i in range(4):
    widen(s, "输入标题", w=2.5, h=0.5, nth=i)
    widen(s, "Add your title here pleases", w=2.5, h=1.15, nth=i)
for i, (t_, d_) in enumerate(SEC):
    set_by_orig(s, "输入标题", t_, nth=0, size=11, bold=True)
    set_by_orig(s, "Add your title here pleases", d_, nth=0, size=8.5)
for i, num in enumerate(["01", "02", "03", "04"]):
    set_by_orig(s, ["1994", "07", "04", "356"][i], num, size=20, bold=True)
built.append(s)

# 10 安全验证：两个真实场景 ─────────────────────────────
page("10 安全验证两场景  <- s11")
s = duplicate(11)
set_chip(s, "安全验证：两个真实场景", "VALIDATION")
sq = find_shape(s, "矩形 2")
if sq is not None:
    add_text(s, 0.95, 1.4, 3.6, 2.6, [
        ("结论", {"size": 14, "bold": True}),
        "",
        "两个场景均已完成真实工具调用测试。",
        "",
        "Agent 无法在无凭证时触发开锁动作，",
        "安全判定始终由 Jetson 本地完成。",
    ], size=10.5, spacing=1.4, after=4)
set_by_orig(s, "单击此处输入标题文本", "有效凭证 +「请开门」", nth=0, size=12,
            bold=True)
set_by_orig(s, "I love you more than",
            "工具调用：current_auth_context → request_unlock；干运行阻止真实动作。",
            nth=0, size=9)
set_by_orig(s, "单击此处输入标题文本", "无效凭证 +「请开门」", nth=0, size=12,
            bold=True)
set_by_orig(s, "I love you more than",
            "工具调用：仅 current_auth_context；不调用开锁工具。",
            nth=0, size=9)
built.append(s)

# 11 异常事件闭环 ──────────────────────────────────────
page("11 异常事件闭环  <- s13")
s = duplicate(13)
tb = find_shape(s, "TextBox 11")      # 原框 2.6x0.48 装不下闭环说明，加宽不越界
if tb is not None:
    tb.width = Inches(3.25)
    tb.height = Inches(1.45)
set_by_orig(s, "我们的产品", "异常事件闭环", size=14, bold=True)
set_by_orig(s, "Ipsum dolor sit amet",
            "触发条件：① 红外持续检测到人员但超过 3 秒仍无法识别人脸；② 完整认证进入融合阶段后失败（活体 / 声纹 / 身份一致性未通过）。",
            size=9.5)
set_by_orig(s, "The example text goes here",
            "链路：认证失败 → localhost:8790/event → 原子写入 latest_event.json → SCP 同步 → WorkBuddy 定时读取 → Bot 通知 → 标记 processed=true。"
            "同一次停留只写一次，离开超过 5 秒重新布防；服务异常只记日志并重试，不改变认证或开锁结果。",
            size=9.5)
built.append(s)

# 12 部署、端口与进程治理 ────────────────────────────────
page("12 部署、端口与进程治理  <- s18")
s = duplicate(18)
set_chip(s, "部署、端口与进程治理", "DEPLOYMENT")
chart = find_shape(s, "Chart 4")
if chart is not None:
    chart._element.getparent().remove(chart._element)
add_text(s, 0.82, 1.55, 4.26, 0.4, "端口分工", size=13, bold=True)
add_text(s, 0.82, 2.0, 4.26, 2.3, [
    "3300   FastGPT · Windows/WSL",
    "8787   门锁工具网关 · Jetson",
    "8790   异常事件服务 · 本机",
    "",
    "配置入口：FASTGPT_PORT、",
    "LOCK_TOOL_GATEWAY_PORT、",
    "LOCK_EVENT_SERVICE_PORT",
], size=10.5, spacing=1.5, after=5)
DEP = [("3300", "FastGPT 工作流与 DeepSeek 调用入口"),
       ("8787", "门锁工具网关：两个受控工具"),
       ("8790", "异常事件服务：仅监听本机"),
       ("一键管理", "run_smart_lock.sh start|stop"),
       ("进程监督", "异常退出自动拉起，清理 PID")]
for i, (t_, d_) in enumerate(DEP):
    set_by_orig(s, ["One Category", "Two Category", "Three Category",
                    "Four Category", "Five Category"][i], t_, size=11,
                bold=True)
    set_by_orig(s, "When an unknown printer took a galley of type", d_,
                nth=0, size=8.5)
built.append(s)

# 13 测试与验证结果 ────────────────────────────────────
page("13 测试与验证结果  <- s21 + 表格")
s = duplicate(21)
set_chip(s, "测试与验证结果", "TEST RESULTS")
TBL13 = [["测试项", "结果"],
         ["原硬件五步认证流程", "开发板已跑通"],
         ["FastGPT Web localhost:3300", "返回 200"],
         ["Agent 有效 / 无效凭证调用", "两场景均通过"],
         ["工具网关生命周期", "含 crash / recovery 通过"],
         ["异常服务生命周期", "含 crash / recovery 通过"],
         ["Pipecat 单元与文件语音链路", "通过"],
         ["最新异常触发与 aplay 改动", "待板端复测"]]
add_table(s, 1.04, 1.15, 5.77, 3.7, TBL13, [3.5, 2.27], size=10)
set_by_orig(s, "Mountain Due", [
    "资源观察", "", "整机约 5.5 GiB", "剩余约 1.6 GiB", "",
    "适用 Orin NX 级设备", "不等同于 Nano 4GB",
], size=8.5, bold=True)
built.append(s)

# 14 关键工程问题与解法 ─────────────────────────────────
page("14 关键工程问题与解法  <- s21 + 表格")
s = duplicate(21)
set_chip(s, "关键工程问题与解法", "ENGINEERING")
TBL14 = [["问题", "解决方案"],
         ["红外持续为 true 重复播报", "上升沿触发 + 状态保持 + 离开重置"],
         ["模型首次使用延迟高", "GUI 与 FunASR 启动时预加载"],
         ["GUI 声纹与 Agent 抢麦", "Agent 先等待硬件凭证再开麦"],
         ["Agent 可能绕过认证", "一次性凭证 + 本地工具网关"],
         ["FastGPT 无法访问 Jetson", "SSH 本地 / 反向隧道 3300·8787"],
         ["服务异常退出污染端口", "PID 监督与信号清理"],
         ["异常情况没有记录", "独立事件服务 + 原子写入"],
         ["音频路径不统一", "输入 PortAudio，输出 ALSA aplay"]]
add_table(s, 1.04, 1.15, 5.77, 3.7, TBL14, [2.75, 3.02], size=9.5)
set_by_orig(s, "Mountain Due", [
    "工程原则", "", "认证归硬件", "交互归 Agent", "执行归网关", "",
    "问题不改变架构结论",
], size=8.5, bold=True)
built.append(s)

# 15 当前边界与风险 ────────────────────────────────────
page("15 当前边界与风险  <- s8")
s = duplicate(8)
set_chip(s, "验证与规划", "VERIFICATION")
t = find_shape(s, "矩形 106")
if t is not None:
    t.width = Inches(4.2)
    t.left = Inches(5.0 - 2.1)
set_by_orig(s, "输入标题", "当前边界与风险", nth=0, size=15, bold=True)
set_by_orig(s, "We are more powerful when we are clear",
            "这些是下一阶段工程工作，不影响已完成的架构验证结论。", size=9.5)
set_by_orig(s, "Clarity Leads to Power.", "", size=9)
RISK = [("板端同步未完成", "开发板 SSH 在密钥交换阶段断开，最新改动尚未同步复测。"),
        ("执行器仍为模拟", "当前用串口风扇模拟开锁，真实继电器需有人看护并可断电。"),
        ("DeepSeek 依赖网络", "断网时硬件认证可运行，agent_confirm 对话开门不可用。"),
        ("稳定性数据缺失", "回声消除、长时间打断与端到端延迟尚无正式测量数据。")]
for i, (t_, d_) in enumerate(RISK):
    set_by_orig(s, "Header here please", t_, nth=0, size=12, bold=True)
    set_by_orig(s, "We are placing great text in this text frames.", d_,
                nth=0, size=9)
built.append(s)

# 16 下一步计划 ────────────────────────────────────────
page("16 下一步计划  <- s20")
s = duplicate(20)
set_chip(s, "下一步计划", "NEXT STEPS")
set_by_orig(s, "Creepiest god air fish land darkness air saying.",
            "短期完成板端验收，中期形成可演示闭环，长期走向产品化。", size=9.5)
ROWS16 = [
    ("Also he his male air bring is signs first air bring is signs.",
     "短期：恢复 SSH 同步，复测板端事件"),
    ("Creepiest god air fish land darkness air saying. Also he his male air.",
     "中期：接入通知渠道，走通闭环演示"),
    ("Also he his male air bring is signs first.",
     "长期：SQLite 队列、真实继电器压测"),
]
for i, (orig, new) in enumerate(ROWS16):
    set_by_orig(s, orig, new, size=9)
STAGE = [("短期", "板端验收复测"), ("中期", "通知闭环演示"), ("长期", "产品化与压测")]
for i, (num, lab) in enumerate(STAGE):
    set_by_orig(s, ["Also his", "Also his male", "His male"][i], num,
                size=17, bold=True)
    set_by_orig(s, ["Bring is signs first Creepiest god air.",
                    "Air bring is signs first Creepiest god air.",
                    "Signs first Creepiest god air."][i], lab, size=9)
built.append(s)

# 17 感谢 ─────────────────────────────────────────────
page("17 感谢  <- s22")
s = duplicate(22)
set_by_orig(s, "感谢观看", "感谢观看", size=30, bold=True)
set_by_orig(s, "商务通用工作总结ppt模板",
            "硬件认证 · Agent 交互 · 网关执行", size=13)
set_by_orig(s, "Business Work Summary Report Template",
            "THANKS FOR WATCHING", size=10)
set_by_orig(s, "汇报人", "MVP 初稿演示", size=9)
set_by_orig(s, "XXXX", "2026.08", size=9)
built.append(s)

# ─────────────────────────── 收尾 ───────────────────────────
n_new = len(built)
delete_slides(list(range(0, len(prs.slides) - n_new)))

prs.core_properties.title = "基于 Jetson 的多模态智能门锁与语音 Agent"
prs.core_properties.subject = "MVP 初稿演示"
prs.save(OUT)

print(f"\n{'=' * 60}")
print(f"saved: {OUT}")
print(f"slides: {len(prs.slides)}  (built={n_new})")
print(f"canvas: {Emu(prs.slide_width).inches} x {Emu(prs.slide_height).inches}")
