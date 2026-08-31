# -*- coding: utf-8 -*-
"""初稿交付校验：残留占位文案 / 动画继承 / 图片 / 溢出风险。"""
import math
from pptx import Presentation
from pptx.util import Emu

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
OUT = r"C:\Users\hoyo\Desktop\lock\docs\智能门锁_语音Agent_初稿.pptx"

prs = Presentation(OUT)
SW = Emu(prs.slide_width).inches
SH = Emu(prs.slide_height).inches

LEFTOVER = [
    "极简淡蓝", "商务通用工作总结", "Business Work Summary", "输入标题",
    "Add your title here", "请输入你的题目", "请在此处添加详细描述",
    "Header here please", "We are placing great text", "Lorem", "lorem",
    "One Category", "Two Category", "Three Category", "Four Category",
    "Five Category", "When an unknown printer", "单击此处输入标题文本",
    "I love you more than", "我们的产品", "The example text goes",
    "Ipsum dolor sit amet", "Mountain Due", "Dummy text ever since",
    "Client :", "Also he his male", "Also his", "His male",
    "Bring is signs first", "Creepiest god air", "Air saying fruitful",
    "Lorem Ipsum is simply", "Lesser first day", "PLEASE ENTER",
    "汇报人", "XXXX", "1994", "CEO", "MANAGER", "John Marteen",
    "We are more powerful", "Clarity Leads to Power", "Fruitful were",
    "Hath waters divide", "Lesser first day kind", "关于我们", "ABOUT US",
    "标题数字等都可以通过", "Ipsum dolor", "The example text",
]

print(f"slides={len(prs.slides)}  canvas={SW}x{SH}")


def walk_text(shapes, out):
    for sh in shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            out.append((sh, sh.text_frame.text))
        if sh.shape_type == 6:
            try:
                walk_text(sh.shapes, out)
            except Exception:
                pass
        if sh.shape_type == 19:  # table
            for row in sh.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        out.append((sh, cell.text))


def est_lines(sh, txt):
    """粗估行数：按首个 run 字号与框宽估算。"""
    try:
        w = Emu(sh.width).inches
    except Exception:
        return 0
    size = 12.0
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size:
                size = r.font.size.pt
                break
        break
    # 中文按 1.0 em/字，英文按 0.55 em/字
    units = 0.0
    for ch in txt:
        units += 1.0 if ord(ch) > 0x2E80 else 0.55
    px_per_line = max(w - 0.12, 0.3) * 72.0
    return units * size / px_per_line


total_anim = 0
problems = []

for i, s in enumerate(prs.slides):
    texts = []
    walk_text(s.shapes, texts)
    blob = "\n".join(t for _, t in texts)

    hits = [w for w in LEFTOVER if w in blob]
    if hits:
        problems.append(f"s{i:02d} 残留占位文案: {hits}")

    timing = s.element.find(f"{{{NS_P}}}timing")
    n_anim = 0
    if timing is not None:
        for el in timing.iter():
            if el.tag.split("}")[-1] in ("animEffect", "animMotion",
                                         "animRot", "animScale"):
                n_anim += 1
    total_anim += n_anim

    n_pic = sum(1 for sh in s.shapes if sh.shape_type == 13)
    n_tbl = sum(1 for sh in s.shapes if sh.shape_type == 19)

    # 溢出粗估
    ovf = []
    for sh, txt in texts:
        if sh.shape_type == 19:
            continue
        try:
            h = Emu(sh.height).inches
        except Exception:
            continue
        size = 12.0
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    size = r.font.size.pt
                    break
            break
        lines = est_lines(sh, txt)
        need = lines * size * 1.45 / 72.0
        if need > h * 1.55 and lines > 2:
            ovf.append(f"{sh.name[:14]}(需{need:.2f}in/框{h:.2f}in)")
    if ovf:
        problems.append(f"s{i:02d} 可能溢出: {ovf[:3]}")

    # 越界
    for sh in s.shapes:
        try:
            l, t = Emu(sh.left).inches, Emu(sh.top).inches
            w, h = Emu(sh.width).inches, Emu(sh.height).inches
        except Exception:
            continue
        if l < -0.35 or t < -0.35 or l + w > SW + 0.35 or t + h > SH + 0.35:
            if sh.shape_type != 6:
                problems.append(
                    f"s{i:02d} 越界: {sh.name[:16]} "
                    f"({l:.2f},{t:.2f},{w:.2f},{h:.2f})")

    title = ""
    for sh, txt in texts:
        if txt.strip():
            title = txt.strip().replace("\n", " / ")[:44]
            break
    print(f"  s{i:02d} anim={n_anim:<3} pic={n_pic} tbl={n_tbl} "
          f"txt={len(texts):<3} | {title}")

print(f"\ntotal animation effects = {total_anim}")

# 图片清单
print("\n图片引用：")
broken = 0
for i, s in enumerate(prs.slides):
    rels = {r for r in s.part.rels}
    for sh in s.shapes:
        if sh.shape_type == 13:
            b = sh._element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip")
            rid = b.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed") if b is not None else None
            if rid not in rels:
                broken += 1
                print(f"  s{i:02d} BROKEN IMAGE {sh.name} rid={rid}")
                continue
            print(f"  s{i:02d} {sh.name[:18]} "
                  f"{sh.image.filename} {sh.image.size}")
print(f"\n断裂图片引用: {broken}")

print("\n问题：")
if problems:
    for p in problems:
        print("  !", p)
else:
    print("  无")
