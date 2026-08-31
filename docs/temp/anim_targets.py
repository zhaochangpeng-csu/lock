# -*- coding: utf-8 -*-
"""列出每页动画目标形状 id 与效果类型，确定哪些形状必须保留。"""
from collections import Counter
from pptx import Presentation

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

PATH = r"C:\Users\hoyo\Desktop\lock\14 (1).pptx"
prs = Presentation(PATH)

for i, s in enumerate(prs.slides):
    timing = s.element.find(f"{{{NS_P}}}timing")
    if timing is None:
        continue
    ids = Counter()
    for spTgt in timing.iter(f"{{{NS_P}}}spTgt"):
        sid = spTgt.get("spid")
        ids[sid] += 1
    # 效果类型
    kinds = Counter()
    for el in timing.iter():
        tag = el.tag.split("}")[-1]
        if tag in ("animEffect", "animMotion", "animRot", "animScale",
                   "animClr", "set", "animate"):
            kinds[tag] += 1
    # 形状 id -> 名称
    id2name = {}
    for sh in s.shapes:
        id2name[str(sh.shape_id)] = sh.name[:20]
        if sh.shape_type == 6:
            try:
                for c in sh.shapes:
                    id2name[str(c.shape_id)] = "  " + c.name[:18]
            except Exception:
                pass
    if ids:
        tg = ", ".join(f"{sid}({id2name.get(sid,'?')})x{n}"
                       for sid, n in ids.most_common(6))
        print(f"s{i:<2} fx={dict(kinds)}")
        print(f"    targets: {tg}")
